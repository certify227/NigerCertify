#!/usr/bin/env python3
"""PowerPoint — Vidéo 1 : IP, masque, passerelle (1 h) — ABOU SAYABOU."""

from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(__file__).resolve().parent / "Video1_IP_Masque_Passerelle_ABOU_SAYABOU.pptx"

# Couleurs
NAVY = RGBColor(0x0D, 0x3B, 0x66)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
LIGHT = RGBColor(0xE8, 0xF1, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
RED = RGBColor(0xC6, 0x28, 0x28)

# 16:9
W, H = Inches(13.333), Inches(7.5)


def set_run(run, size=20, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return shape


def add_top_bar(slide, title_right="Cours réseau — 1 h"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  ABOU SAYABOU  •  Réseau informatique"
    set_run(run, 12, True, WHITE)
    # right label via second shape
    lab = slide.shapes.add_textbox(Inches(8.5), Inches(0.12), Inches(4.5), Inches(0.35))
    p = lab.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = title_right + "  "
    set_run(run, 12, False, LIGHT)


def add_footer(slide, page, total):
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), W, Inches(0.35))
    foot.fill.solid()
    foot.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    foot.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.18), Inches(10), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Vidéo 1 — IP, masque, passerelle  |  Débutant → Intermédiaire"
    set_run(run, 10, False, GRAY)
    tb2 = slide.shapes.add_textbox(Inches(11.5), Inches(7.18), Inches(1.5), Inches(0.3))
    p = tb2.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page}/{total}"
    set_run(run, 10, True, GRAY)


def add_title(slide, text, top=0.7):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, 32, True, NAVY)
    # accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top + 0.65), Inches(1.8), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_subtitle(slide, text, top=1.45):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, 16, False, GRAY)


def add_bullets(slide, items, left=0.6, top=2.0, width=12, size=20):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "▸  " + item
        set_run(run, size, False, DARK)


def add_card(slide, left, top, width, height, title, body, fill=LIGHT, title_color=NAVY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run(run, 16, True, title_color)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    run = p2.add_run()
    run.text = body
    set_run(run, 13, False, DARK)
    return card


def add_timing(slide, text):
    """Badge timing en haut à droite du contenu."""
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.6), Inches(0.7), Inches(2.3), Inches(0.4)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    set_run(run, 12, True, WHITE)


def blank_content_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slides_meta = []  # for page numbers later — we'll add footers at end by re-opening... easier to count
    TOTAL = 22  # planned

    # ========== 1 COVER ==========
    s = blank_content_slide(prs)
    add_bg(s, NAVY)
    # accent strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.8), W, Inches(1.7))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(0x0A, 0x2F, 0x52)
    strip.line.fill.background()

    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "COURS RÉSEAU  •  VIDÉO 1"
    set_run(run, 18, True, ACCENT)

    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.5), Inches(2))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "IP, masque, passerelle :"
    set_run(run, 40, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    run = p2.add_run()
    run.text = "le tuto qui règle la confusion"
    set_run(run, 36, True, LIGHT)

    tb = s.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Durée 1 heure  •  Niveau débutant → intermédiaire  •  Théorie + lab + exercices"
    set_run(run, 16, False, RGBColor(0xB0, 0xC4, 0xDE))

    tb = s.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "ABOU SAYABOU"
    set_run(run, 22, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    run = p2.add_run()
    run.text = "Formateur cybersécurité & réseaux"
    set_run(run, 14, False, LIGHT)

    # ========== 2 OBJECTIF ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Objectif de la vidéo")
    add_timing(s, "0:00 – 1:30")
    add_subtitle(s, "À la fin de cette heure, tu sauras…")
    add_card(s, 0.5, 2.1, 3.9, 3.8, "1. Lire une IP",
             "Comprendre les 4 octets, le rôle de chaque partie, et ce qu’est une adresse privée vs publique.",
             LIGHT)
    add_card(s, 4.7, 2.1, 3.9, 3.8, "2. Utiliser le masque / CIDR",
             "Calculer réseau, broadcast et plage d’hôtes pour /24, /16, /26… sans calculatrice magique.",
             LIGHT)
    add_card(s, 8.9, 2.1, 3.9, 3.8, "3. Diagnostiquer",
             "Savoir si le problème vient de l’IP, du masque, de la passerelle ou du DNS (4 ping).",
             LIGHT)
    add_footer(s, 2, TOTAL)

    # ========== 3 PLAN ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Plan de la session (1 h)")
    add_timing(s, "Chapitres")
    plan = [
        ("01", "1:30", "Analogie ville — adresse, quartier, sortie"),
        ("02", "6:00", "IPv4 et binaire (juste ce qu’il faut)"),
        ("03", "14:00", "Masque & CIDR (/8, /16, /24, /25, /30)"),
        ("04", "24:00", "Réseau • Host • Broadcast"),
        ("05", "34:00", "Passerelle (gateway)"),
        ("06", "40:00", "DHCP vs IP fixe"),
        ("07", "45:00", "Lab Windows / Linux"),
        ("08", "52:00", "5 exercices + récap"),
    ]
    y = 1.5
    for num, t, label in plan:
        # num box
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(0.7), Inches(0.48))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = box.text_frame.paragraphs[0].add_run()
        run.text = num
        set_run(run, 14, True, WHITE)

        tbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.35), Inches(y), Inches(1.3), Inches(0.48))
        tbox.fill.solid()
        tbox.fill.fore_color.rgb = ACCENT
        tbox.line.fill.background()
        tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tbox.text_frame.paragraphs[0].add_run()
        run.text = t
        set_run(run, 13, True, WHITE)

        tb = s.shapes.add_textbox(Inches(2.9), Inches(y + 0.05), Inches(9.5), Inches(0.4))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = label
        set_run(run, 18, False, DARK)
        y += 0.58
    add_footer(s, 3, TOTAL)

    # ========== 4 ANALOGIE ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Analogie : le réseau comme une ville")
    add_timing(s, "1:30 – 6:00")
    add_card(s, 0.5, 1.7, 4.0, 2.2, "🏠 Adresse IP",
             "Le numéro de ta maison.\nEx. : 192.168.1.45", LIGHT)
    add_card(s, 4.7, 1.7, 4.0, 2.2, "🗺️ Masque / quartier",
             "Qui est dans ta rue (réseau local).\nEx. : /24 = tout le 192.168.1.x", LIGHT)
    add_card(s, 8.9, 1.7, 4.0, 2.2, "🚪 Passerelle",
             "La sortie de la ville (box / routeur).\nEx. : 192.168.1.1", LIGHT)
    add_bullets(s, [
        "Sans adresse → tu n’existes pas sur le réseau",
        "Sans masque → tu ne sais pas qui est « local »",
        "Sans passerelle → tu restes bloqué dans ton quartier (pas d’Internet)",
    ], top=4.2, size=18)
    add_footer(s, 4, TOTAL)

    # ========== 5 SCHEMA ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Schéma de base à retenir")
    add_timing(s, "Schéma")
    # boxes PC / SW / GW / NET
    nodes = [
        (0.8, "PC1\n192.168.1.10", BLUE),
        (3.5, "PC2\n192.168.1.20", BLUE),
        (6.2, "Switch", GRAY),
        (8.9, "Passerelle\n192.168.1.1", ORANGE),
        (11.2, "Internet", GREEN),
    ]
    for x, text, col in nodes:
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(1.9), Inches(1.5))
        sh.fill.solid()
        sh.fill.fore_color.rgb = col
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = text
        set_run(run, 14, True, WHITE)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(12.3), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = [
        "Réseau local (LAN) : 192.168.1.0/24",
        "Masque : 255.255.255.0",
        "Tout trafic hors 192.168.1.x → envoyé à la passerelle 192.168.1.1",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "▸  " + line
        set_run(run, 18, False, DARK)
    add_footer(s, 5, TOTAL)

    # ========== 6 IPV4 ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "IPv4 : structure d’une adresse")
    add_timing(s, "6:00 – 14:00")
    add_subtitle(s, "4 octets • 32 bits • valeur 0–255 chacun")

    # big IP visual
    octets = [("192", "octet 1"), ("168", "octet 2"), ("1", "octet 3"), ("45", "octet 4")]
    x = 1.2
    for val, lab in octets:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(2.2), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        tf = box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = val
        set_run(run, 36, True, WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run = p2.add_run()
        run.text = lab
        set_run(run, 12, False, LIGHT)
        if x < 8:
            dot = s.shapes.add_textbox(Inches(x + 2.15), Inches(2.6), Inches(0.4), Inches(0.6))
            run = dot.text_frame.paragraphs[0].add_run()
            run.text = "."
            set_run(run, 36, True, ACCENT)
        x += 2.7

    add_bullets(s, [
        "Exemple : 192.168.1.45  →  binaire sur 32 bits",
        "Adresses privées (RFC1918) : 10.0.0.0/8 • 172.16.0.0/12 • 192.168.0.0/16",
        "Les classes A/B/C = historique → aujourd’hui on parle CIDR",
    ], top=4.3, size=18)
    add_footer(s, 6, TOTAL)

    # ========== 7 BINAIRE ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Binaire : juste ce qu’il faut")
    add_timing(s, "Astuce")
    add_subtitle(s, "Un octet = 8 bits. Poids : 128 64 32 16 8 4 2 1")

    # table as cards
    headers = ["128", "64", "32", "16", "8", "4", "2", "1"]
    bits = ["1", "1", "0", "0", "0", "0", "0", "0"]  # 192
    x = 1.0
    for h, b in zip(headers, bits):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(1.3), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY if b == "1" else RGBColor(0x90, 0xA4, 0xAE)
        box.line.fill.background()
        tf = box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = h
        set_run(run, 14, True, WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        run = p2.add_run()
        run.text = b
        set_run(run, 28, True, WHITE)
        x += 1.45

    add_bullets(s, [
        "192 = 128+64 → bits 11000000",
        "255 = tous les bits à 1 → 11111111",
        "0 = tous les bits à 0 → 00000000",
        "Tu n’as PAS besoin d’être un crack du binaire pour utiliser /24 au quotidien",
    ], top=4.4, size=18)
    add_footer(s, 7, TOTAL)

    # ========== 8 MASQUE ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Masque de sous-réseau & CIDR")
    add_timing(s, "14:00 – 24:00")
    add_subtitle(s, "Le masque dit : « quels bits font partie du RÉSEAU ? »")

    rows = [
        ("/8", "255.0.0.0", "~16 millions d’hôtes", "Très grand réseau"),
        ("/16", "255.255.0.0", "65 534 hôtes", "Ex. 172.16.0.0/16"),
        ("/24", "255.255.255.0", "254 hôtes", "Le plus courant (LAN)"),
        ("/25", "255.255.255.128", "126 hôtes", "Découpe un /24 en 2"),
        ("/30", "255.255.255.252", "2 hôtes", "Lien point-à-point"),
    ]
    y = 1.9
    # header
    for i, h in enumerate(["CIDR", "Masque", "Hôtes utiles", "Usage"]):
        xs = [0.5, 2.2, 5.5, 8.8][i]
        ws = [1.5, 3.1, 3.1, 3.8][i]
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(xs), Inches(y), Inches(ws), Inches(0.45))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = box.text_frame.paragraphs[0].add_run()
        run.text = h
        set_run(run, 14, True, WHITE)
    y = 2.4
    for i, row in enumerate(rows):
        bg = LIGHT if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            xs = [0.5, 2.2, 5.5, 8.8][j]
            ws = [1.5, 3.1, 3.1, 3.8][j]
            box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(xs), Inches(y), Inches(ws), Inches(0.5))
            box.fill.solid()
            box.fill.fore_color.rgb = bg
            box.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
            box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = box.text_frame.paragraphs[0].add_run()
            run.text = val
            set_run(run, 14, j == 0, DARK)
        y += 0.5

    tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.2))
    tip.fill.solid()
    tip.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    tip.line.color.rgb = ORANGE
    tf = tip.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = "Formule hôtes = 2^(32 − préfixe) − 2"
    set_run(run, 16, True, ORANGE)
    p2 = tf.add_paragraph()
    run = p2.add_run()
    run.text = "On retire 2 : adresse réseau + adresse broadcast. Ex. /24 → 2^8 − 2 = 254."
    set_run(run, 14, False, DARK)
    add_footer(s, 8, TOTAL)

    # ========== 9 METHODE ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Méthode : réseau / hosts / broadcast")
    add_timing(s, "24:00 – 34:00")
    add_subtitle(s, "Exemple : 192.168.10.45/24")

    steps = [
        ("1", "Écrire le masque", " /24 → 255.255.255.0"),
        ("2", "Trouver le réseau", " Bits hôtes à 0 → 192.168.10.0"),
        ("3", "Trouver le broadcast", " Bits hôtes à 1 → 192.168.10.255"),
        ("4", "Plage utilisable", " 192.168.10.1  →  192.168.10.254"),
    ]
    y = 1.9
    for num, title, detail in steps:
        n = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y), Inches(0.55), Inches(0.55))
        n.fill.solid()
        n.fill.fore_color.rgb = ACCENT
        n.line.fill.background()
        n.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = n.text_frame.paragraphs[0].add_run()
        run.text = num
        set_run(run, 16, True, WHITE)
        tb = s.shapes.add_textbox(Inches(1.4), Inches(y), Inches(11), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title + "  —  "
        set_run(run, 18, True, NAVY)
        run = p.add_run()
        run.text = detail
        set_run(run, 18, False, DARK)
        y += 0.85

    add_footer(s, 9, TOTAL)

    # ========== 10 EXEMPLE /26 ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "2ᵉ exemple : 192.168.10.45/26")
    add_timing(s, "Exercice guidé")
    add_card(s, 0.5, 1.7, 6.0, 4.5, "Calcul",
             "/26 = 255.255.255.192\n"
             "Blocs de 64 adresses (2^6)\n\n"
             "Réseaux : .0  |  .64  |  .128  |  .192\n\n"
             "45 est dans le bloc .0 → .63\n\n"
             "Réseau     : 192.168.10.0\n"
             "Broadcast  : 192.168.10.63\n"
             "Hôtes      : .1 → .62\n"
             "Passerelle typique : .1",
             LIGHT)
    add_card(s, 6.8, 1.7, 5.9, 4.5, "À retenir",
             "• Plus le préfixe est GRAND,\n  plus le réseau est PETIT\n\n"
             "• /24 = confort débutant\n• /30 = liens routeurs\n• /16 = campus / labo\n\n"
             "• Toujours vérifier que la\n  passerelle est DANS le même\n  réseau que l’IP de la machine",
             RGBColor(0xE8, 0xF5, 0xE9), GREEN)
    add_footer(s, 10, TOTAL)

    # ========== 11 PASSERELLE ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "La passerelle (default gateway)")
    add_timing(s, "34:00 – 40:00")
    add_bullets(s, [
        "C’est le « prochain saut » (next hop) pour tout trafic hors LAN",
        "En général : l’IP du routeur / de la box sur ton réseau (souvent .1 ou .254)",
        "Elle DOIT être dans le même sous-réseau que ta machine",
        "Symptôme classique : ping local OK, Internet KO → suspecter GW / NAT / WAN",
    ], top=1.7, size=18)

    add_card(s, 0.5, 4.2, 4.0, 2.2, "Ping voisin", "OK → LAN / switch OK", RGBColor(0xE8, 0xF5, 0xE9), GREEN)
    add_card(s, 4.7, 4.2, 4.0, 2.2, "Ping passerelle", "OK → ta config L3 locale OK", RGBColor(0xFF, 0xF8, 0xE1), ORANGE)
    add_card(s, 8.9, 4.2, 3.9, 2.2, "Ping 8.8.8.8", "OK → routage Internet OK", LIGHT, NAVY)
    add_footer(s, 11, TOTAL)

    # ========== 12 DHCP ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "DHCP vs adresse IP fixe")
    add_timing(s, "40:00 – 45:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "DHCP (automatique)",
             "• La box / le serveur DHCP donne :\n"
             "  IP + masque + passerelle + DNS\n\n"
             "• Idéal pour PC, téléphones, invités\n"
             "• Bail = durée de validité\n"
             "• Réservation DHCP = « fixe soft »\n"
             "  (même IP selon MAC)\n\n"
             "• APIPA 169.254.x.x = pas de DHCP !",
             LIGHT)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "IP fixe (statique)",
             "• Tu configures tout à la main\n\n"
             "• Idéal pour :\n"
             "  serveurs, imprimantes, NAS,\n"
             "  caméras, passerelles\n\n"
             "• Documente : IP / masque / GW / DNS\n"
             "• Évite les conflits (hors plage DHCP)\n\n"
             "• Erreur fréquente : GW oubliée\n"
             "  ou DNS oublié",
             RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
    add_footer(s, 12, TOTAL)

    # ========== 13 LAB ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Lab : lire sa config")
    add_timing(s, "45:00 – 52:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Windows",
             "ipconfig /all\n\n"
             "• Adresse IPv4\n"
             "• Masque de sous-réseau\n"
             "• Passerelle par défaut\n"
             "• Serveurs DNS\n"
             "• DHCP activé ? Bail ?\n\n"
             "ipconfig /release\n"
             "ipconfig /renew",
             LIGHT, NAVY)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Linux",
             "ip a\n"
             "ip route\n"
             "cat /etc/resolv.conf\n\n"
             "• inet = IP + CIDR (/24)\n"
             "• default via = passerelle\n"
             "• nameserver = DNS\n\n"
             "nmcli device show\n"
             "(NetworkManager)",
             RGBColor(0xE8, 0xF5, 0xE9), GREEN)
    add_footer(s, 13, TOTAL)

    # ========== 14 4 PINGS ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Les 4 ping du diagnostic")
    add_timing(s, "Méthode")
    tests = [
        ("1", "127.0.0.1", "Pile TCP/IP locale", "Si KO → OS / stack réseau"),
        ("2", "Passerelle", "LAN + config L3", "Si KO → câble, IP, VLAN, GW"),
        ("3", "8.8.8.8", "Sortie Internet", "Si KO → NAT / FAI / firewall"),
        ("4", "google.com", "Résolution DNS", "Si KO alors que 3 OK → DNS"),
    ]
    y = 1.7
    for num, target, meaning, tip in tests:
        n = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(0.6), Inches(1.0))
        n.fill.solid()
        n.fill.fore_color.rgb = NAVY
        n.line.fill.background()
        n.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = n.text_frame.paragraphs[0].add_run()
        run.text = "\n" + num
        set_run(run, 20, True, WHITE)

        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(y), Inches(11.4), Inches(1.0))
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT
        c.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"ping {target}   —   {meaning}"
        set_run(run, 16, True, NAVY)
        p2 = tf.add_paragraph()
        run = p2.add_run()
        run.text = tip
        set_run(run, 14, False, DARK)
        y += 1.15
    add_footer(s, 14, TOTAL)

    # ========== 15 EXERCICES ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "5 exercices (à faire avec moi)")
    add_timing(s, "52:00 – 58:00")
    exs = [
        "1.  IP 10.0.0.50/24  →  réseau ? broadcast ? plage ?",
        "2.  IP 172.16.5.10/16  →  réseau ? broadcast ?",
        "3.  192.168.1.10/24 et 192.168.2.10/24 : même réseau ?",
        "4.  PC = 192.168.1.50/24  GW = 192.168.0.1  →  problème ?",
        "5.  Combien d’hôtes dans un /27 ?",
    ]
    add_bullets(s, exs, top=1.8, size=20)
    tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.0))
    tip.fill.solid()
    tip.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    tip.line.color.rgb = ACCENT
    run = tip.text_frame.paragraphs[0].add_run()
    run.text = "Pause écran : laisse 20–30 secondes par exercice avant le corrigé."
    set_run(run, 16, True, NAVY)
    add_footer(s, 15, TOTAL)

    # ========== 16 CORRIGE ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Corrigé des exercices")
    add_timing(s, "Corrigé")
    corrects = [
        "1.  Réseau 10.0.0.0  •  Broadcast 10.0.0.255  •  Hôtes .1–.254",
        "2.  Réseau 172.16.0.0  •  Broadcast 172.16.255.255",
        "3.  NON — 192.168.1.0/24 ≠ 192.168.2.0/24 (il faut un routeur)",
        "4.  OUI problème — GW hors du réseau du PC (mauvais sous-réseau)",
        "5.  /27 → 2^(32-27) − 2 = 30 hôtes",
    ]
    add_bullets(s, corrects, top=1.8, size=19)
    add_footer(s, 16, TOTAL)

    # ========== 17 ERREURS ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Erreurs classiques (à éviter)")
    add_timing(s, "58:00")
    errs = [
        ("Masque faux", "Deux PC « se voient » mal ou pas du tout"),
        ("GW d’un autre réseau", "LAN OK, Internet impossible"),
        ("Conflit d’IP", "Coupures aléatoires, APIPA"),
        ("DNS oublié (IP fixe)", "Ping IP OK, sites web KO"),
        ("169.254.x.x", "Pas de réponse DHCP — câble / serveur DHCP"),
    ]
    y = 1.7
    for title, detail in errs:
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(0.85))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
        c.line.color.rgb = RED
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "✖  " + title + "  →  "
        set_run(run, 16, True, RED)
        run = p.add_run()
        run.text = detail
        set_run(run, 16, False, DARK)
        y += 0.95
    add_footer(s, 17, TOTAL)

    # ========== 18 RECAP ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Récap — 4 phrases à retenir")
    add_timing(s, "58:00 – 59:30")
    recaps = [
        ("IP", "Identifie la machine sur le réseau"),
        ("Masque / CIDR", "Délimite qui est « local »"),
        ("Passerelle", "Sortie pour tout le reste (Internet)"),
        ("4 ping", "Localhost → GW → IP publique → nom DNS"),
    ]
    y = 1.8
    for a, b in recaps:
        left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(3.5), Inches(0.95))
        left.fill.solid()
        left.fill.fore_color.rgb = NAVY
        left.line.fill.background()
        left.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = left.text_frame.paragraphs[0].add_run()
        run.text = "\n" + a
        set_run(run, 18, True, WHITE)
        right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(y), Inches(8.6), Inches(0.95))
        right.fill.solid()
        right.fill.fore_color.rgb = LIGHT
        right.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        run = right.text_frame.paragraphs[0].add_run()
        run.text = "  " + b
        set_run(run, 18, False, DARK)
        y += 1.15
    add_footer(s, 18, TOTAL)

    # ========== 19 CHEAT ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Aide-mémoire (screenshot)")
    add_timing(s, "Fiche")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Commandes",
             "Windows\n"
             "  ipconfig /all\n"
             "  ping 192.168.1.1\n"
             "  nslookup google.com\n\n"
             "Linux\n"
             "  ip a\n"
             "  ip route\n"
             "  ping -c 4 8.8.8.8\n"
             "  dig google.com",
             LIGHT)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "CIDR express",
             "/8   →  255.0.0.0\n"
             "/16  →  255.255.0.0\n"
             "/24  →  255.255.255.0\n"
             "/25  →  255.255.255.128\n"
             "/26  →  255.255.255.192\n"
             "/27  →  255.255.255.224\n"
             "/30  →  255.255.255.252\n\n"
             "Hôtes = 2^(32-n) − 2",
             RGBColor(0xE8, 0xF5, 0xE9), GREEN)
    add_footer(s, 19, TOTAL)

    # ========== 20 PROCHAINE ==========
    s = blank_content_slide(prs)
    add_top_bar(s)
    add_title(s, "Prochaine vidéo + devoir")
    add_timing(s, "59:30 – 60:00")
    add_card(s, 0.5, 1.8, 6.0, 3.5, "Vidéo 2 (semaine prochaine)",
             "VLAN expliqué simplement\n+ lab complet\n\n3 VLAN • Access • Trunk\nInter-VLAN routing",
             LIGHT, NAVY)
    add_card(s, 6.8, 1.8, 5.9, 3.5, "Devoir (commentaire)",
             "Calcule pour 10.0.5.80/23 :\n\n• adresse réseau\n• broadcast\n• plage d’hôtes\n\nPoste ta réponse en commentaire !",
             RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
    add_bullets(s, [
        "Abonne-toi pour la playlist « Réseau 1 h »",
        "Like si le masque est enfin clair 👍",
    ], top=5.5, size=18)
    add_footer(s, 20, TOTAL)

    # ========== 21 FIN ==========
    s = blank_content_slide(prs)
    add_bg(s, NAVY)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Merci !"
    set_run(run, 48, True, WHITE)

    tb = s.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.5), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "IP  •  Masque  •  Passerelle  •  4 ping"
    set_run(run, 22, False, LIGHT)

    tb = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "ABOU SAYABOU"
    set_run(run, 24, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run = p2.add_run()
    run.text = "Cours réseau — Vidéo 1 / 3"
    set_run(run, 14, False, ACCENT)

    # ========== 22 NOTES FORMATEUR (optionnel) ==========
    s = blank_content_slide(prs)
    add_top_bar(s, "Notes formateur (ne pas montrer)")
    add_title(s, "Notes de tournage")
    add_bullets(s, [
        "Garder le schéma slide 5 visible en fond pendant le lab",
        "Pause 20–30 s sur chaque exercice avant corrigé",
        "Montrer ipconfig /all en vrai (pas seulement capture)",
        "Répéter : « GW dans le même réseau que l’IP »",
        "CTA commentaire : 10.0.5.80/23 (réponse : réseau 10.0.4.0, BC 10.0.5.255)",
        "Teaser vidéo 2 VLAN à la fin — même charte graphique",
        "Durée cible parlée : ~55–60 min avec lab",
    ], top=1.7, size=17)
    add_footer(s, 22, TOTAL)

    # Fix footers page numbers on slides that have them — already hardcoded
    # Cover and last navy slides don't need footer

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)

    # Métadonnées auteur (retirer python-pptx)
    now = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    buf = BytesIO(OUT.read_bytes())
    out_buf = BytesIO()
    with zipfile.ZipFile(buf, "r") as zin, zipfile.ZipFile(out_buf, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "docProps/core.xml":
                data = f"""<?xml version='1.0' encoding='UTF-8' standalone='yes'?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:dcmitype="http://purl.org/dc/dcmitype/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <dc:title>Vidéo 1 — IP, masque, passerelle (cours 1h)</dc:title>
  <dc:subject>Cours réseau débutant — PowerPoint YouTube</dc:subject>
  <dc:creator>ABOU SAYABOU</dc:creator>
  <cp:keywords>réseau, IP, CIDR, passerelle, YouTube</cp:keywords>
  <dc:description>Support de cours vidéo 1h — IP, masque, passerelle</dc:description>
  <cp:lastModifiedBy>ABOU SAYABOU</cp:lastModifiedBy>
  <cp:revision>1</cp:revision>
  <dcterms:created xsi:type="dcterms:W3CDTF">{now}</dcterms:created>
  <dcterms:modified xsi:type="dcterms:W3CDTF">{now}</dcterms:modified>
  <cp:category>Formation réseau</cp:category>
  <dc:language>fr-FR</dc:language>
</cp:coreProperties>""".encode("utf-8")
            elif item.filename == "docProps/app.xml":
                text = data.decode("utf-8")
                text = text.replace("<Company/>", "<Company>ABOU SAYABOU</Company>")
                text = text.replace("python-pptx", "ABOU SAYABOU")
                data = text.encode("utf-8")
            zout.writestr(item, data)
    OUT.write_bytes(out_buf.getvalue())
    print(f"PowerPoint généré : {OUT} ({OUT.stat().st_size} octets)")


if __name__ == "__main__":
    build()
