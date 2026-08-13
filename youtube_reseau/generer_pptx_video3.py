#!/usr/bin/env python3
"""PowerPoint — Vidéo 3 : Diagnostic panne réseau 7 étapes (1 h) — ABOU SAYABOU."""

from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Video3_Diagnostic_7_Etapes_ABOU_SAYABOU.pptx"

NAVY = RGBColor(0x0D, 0x3B, 0x66)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
LIGHT = RGBColor(0xE8, 0xF1, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
RED = RGBColor(0xC6, 0x28, 0x28)
YELLOW_BG = RGBColor(0xFF, 0xF8, 0xE1)
RED_BG = RGBColor(0xFF, 0xEB, 0xEE)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)

W, H = Inches(13.333), Inches(7.5)
TOTAL = 24
FOOTER = "Vidéo 3 — Diagnostic 7 étapes  |  Débutant → Intermédiaire"


def set_run(run, size=20, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_top_bar(slide, title_right="Cours réseau — 1 h"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    run = bar.text_frame.paragraphs[0].add_run()
    run.text = "  ABOU SAYABOU  •  Réseau informatique"
    set_run(run, 12, True, WHITE)
    lab = slide.shapes.add_textbox(Inches(8.5), Inches(0.12), Inches(4.5), Inches(0.35))
    p = lab.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = title_right + "  "
    set_run(run, 12, False, LIGHT)


def add_footer(slide, page):
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), W, Inches(0.35))
    foot.fill.solid()
    foot.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    foot.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.18), Inches(10), Inches(0.3))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = FOOTER
    set_run(run, 10, False, GRAY)
    tb2 = slide.shapes.add_textbox(Inches(11.5), Inches(7.18), Inches(1.5), Inches(0.3))
    p = tb2.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page}/{TOTAL}"
    set_run(run, 10, False, GRAY)


def add_title(slide, text, top=0.7):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.7))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    set_run(run, 30, True, NAVY)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top + 0.62), Inches(1.8), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_subtitle(slide, text, top=1.42):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.4))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    set_run(run, 16, False, GRAY)


def add_bullets(slide, items, left=0.6, top=2.0, width=12, size=18):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "▸  " + item
        set_run(run, size, False, DARK)


def add_card(slide, left, top, width, height, title, body, fill=LIGHT, title_color=NAVY, size=13):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    run = tf.paragraphs[0].add_run()
    run.text = title
    set_run(run, 16, True, title_color)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    run = p2.add_run()
    run.text = body
    set_run(run, size, False, DARK)
    return card


def add_timing(slide, text):
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.68), Inches(2.4), Inches(0.4)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    set_run(run, 12, True, WHITE)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_code_card(slide, left, top, width, height, title, code):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x32)
    card.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.1)
    run = tf.paragraphs[0].add_run()
    run.text = title
    set_run(run, 13, True, ACCENT, "Calibri")
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    run = p2.add_run()
    run.text = code
    set_run(run, 13, False, WHITE, "Consolas")


def stamp_meta(path: Path):
    now = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    buf = BytesIO(path.read_bytes())
    out_buf = BytesIO()
    with zipfile.ZipFile(buf, "r") as zin, zipfile.ZipFile(
        out_buf, "w", compression=zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "docProps/core.xml":
                data = f"""<?xml version='1.0' encoding='UTF-8' standalone='yes'?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <dc:title>Vidéo 3 — Diagnostic panne réseau 7 étapes (cours 1h)</dc:title>
  <dc:subject>Cours réseau — dépannage OSI, 4 ping, incidents</dc:subject>
  <dc:creator>ABOU SAYABOU</dc:creator>
  <cp:keywords>diagnostic réseau, OSI, ping, DNS, YouTube</cp:keywords>
  <dc:description>Support de cours vidéo 1h — diagnostic 7 étapes</dc:description>
  <cp:lastModifiedBy>ABOU SAYABOU</cp:lastModifiedBy>
  <cp:revision>1</cp:revision>
  <dcterms:created xsi:type="dcterms:W3CDTF">{now}</dcterms:created>
  <dcterms:modified xsi:type="dcterms:W3CDTF">{now}</dcterms:modified>
  <cp:category>Formation réseau</cp:category>
  <dc:language>fr-FR</dc:language>
</cp:coreProperties>""".encode("utf-8")
            elif item.filename == "docProps/app.xml":
                text = data.decode("utf-8")
                text = text.replace("<Company/>", "<Company>ABOU SAYABOU</Company>")
                text = text.replace("python-pptx", "ABOU SAYABOU")
                data = text.encode("utf-8")
            zout.writestr(item, data)
    path.write_bytes(out_buf.getvalue())


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 COVER
    s = blank(prs)
    add_bg(s, NAVY)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.8), W, Inches(1.7))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(0x0A, 0x2F, 0x52)
    strip.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "COURS RÉSEAU  •  VIDÉO 3 / 3"
    set_run(run, 18, True, ACCENT)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.5), Inches(2.4))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Le réseau est down :"
    set_run(run, 36, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    run = p2.add_run()
    run.text = "7 étapes pour diagnostiquer"
    set_run(run, 32, True, LIGHT)
    p3 = tb.text_frame.add_paragraph()
    run = p3.add_run()
    run.text = "n’importe quelle panne"
    set_run(run, 28, True, ACCENT)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(4.85), Inches(11.5), Inches(0.4))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Durée 1 heure  •  Méthode OSI  •  3 incidents  •  4 ping"
    set_run(run, 16, False, RGBColor(0xB0, 0xC4, 0xDE))
    tb = s.shapes.add_textbox(Inches(0.8), Inches(6.15), Inches(11.5), Inches(0.8))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "ABOU SAYABOU"
    set_run(run, 22, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    run = p2.add_run()
    run.text = "Formateur cybersécurité & réseaux"
    set_run(run, 14, False, LIGHT)

    # 2 OBJECTIF
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Objectif de la vidéo")
    add_timing(s, "0:00 – 2:00")
    add_subtitle(s, "Arrêter de pinger au hasard. Appliquer une méthode.")
    add_card(s, 0.5, 2.0, 3.9, 4.2, "1. Cadrer",
             "Poser les bonnes questions avant les commandes : qui, depuis quand, Wi-Fi ou filaire ?",
             LIGHT)
    add_card(s, 4.7, 2.0, 3.9, 4.2, "2. Isoler la couche",
             "Remonter OSI / TCP-IP : physique → lien → IP → DNS → application.",
             LIGHT)
    add_card(s, 8.9, 2.0, 3.9, 4.2, "3. Résoudre 3 cas",
             "Pas de lien • LAN OK Internet KO • Un site marche, l’autre non.",
             LIGHT)
    add_footer(s, 2)

    # 3 PLAN
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Plan de la session (1 h)")
    add_timing(s, "Chapitres")
    plan = [
        ("01", "0:00", "Accroche : ne plus pinger au hasard"),
        ("02", "2:00", "Vue d’ensemble des 7 étapes"),
        ("03", "6:00", "Étapes 1 → 7 (détail + commandes)"),
        ("04", "44:00", "Incident A — pas de lien"),
        ("05", "48:00", "Incident B — LAN OK, Internet KO"),
        ("06", "53:00", "Incident C — un site KO"),
        ("07", "58:00", "Fiche récap + devoir"),
    ]
    y = 1.6
    for num, t, label in plan:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(0.7), Inches(0.55))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = box.text_frame.paragraphs[0].add_run()
        run.text = num
        set_run(run, 14, True, WHITE)
        tbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.35), Inches(y), Inches(1.4), Inches(0.55))
        tbox.fill.solid()
        tbox.fill.fore_color.rgb = ACCENT
        tbox.line.fill.background()
        tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tbox.text_frame.paragraphs[0].add_run()
        run.text = t
        set_run(run, 13, True, WHITE)
        tb = s.shapes.add_textbox(Inches(3.0), Inches(y + 0.08), Inches(9.5), Inches(0.45))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = label
        set_run(run, 18, False, DARK)
        y += 0.68
    add_footer(s, 3)

    # 4 ACCROCHE
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Débutant vs intermédiaire")
    add_timing(s, "0:00 – 2:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Débutant",
             "« Je ping google.com »\n\n"
             "• Si KO → panique\n"
             "• Si OK → « c’est le PC »\n"
             "• Commandes au hasard\n"
             "• Pas de notes\n"
             "• Change 3 choses d’un coup\n\n"
             "Résultat : perte de temps\net faux diagnostics.",
             RED_BG, RED, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Intermédiaire",
             "« Je pose 3 questions »\n\n"
             "• Qui est touché ?\n"
             "• Depuis quand ?\n"
             "• Qu’est-ce qui a changé ?\n"
             "• Puis : 4 ping dans l’ordre\n"
             "• Une couche à la fois\n\n"
             "Résultat : cause trouvée\nen minutes, pas en heures.",
             GREEN_BG, GREEN, 14)
    add_footer(s, 4)

    # 5 VUE 7 ETAPES
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Les 7 étapes (vue d’ensemble)")
    add_timing(s, "2:00 – 6:00")
    steps = [
        ("1", "Cadrer le symptôme"),
        ("2", "Couche 1 — Physique"),
        ("3", "Couche 2 — Lien / VLAN / ARP"),
        ("4", "Couche 3 — IP / masque / GW"),
        ("5", "DNS vs IP publique"),
        ("6", "Ports / firewall / proxy"),
        ("7", "Journaliser et escalader"),
    ]
    y = 1.55
    for num, label in steps:
        n = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y), Inches(0.55), Inches(0.55))
        n.fill.solid()
        n.fill.fore_color.rgb = ACCENT
        n.line.fill.background()
        n.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = n.text_frame.paragraphs[0].add_run()
        run.text = num
        set_run(run, 16, True, WHITE)
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(y), Inches(11.3), Inches(0.55))
        bar.fill.solid()
        bar.fill.fore_color.rgb = LIGHT
        bar.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        run = bar.text_frame.paragraphs[0].add_run()
        run.text = "  " + label
        set_run(run, 18, True, NAVY)
        y += 0.72
    add_footer(s, 5)

    # 6 ETAPE 1
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Étape 1 — Cadrer le symptôme")
    add_timing(s, "6:00 – 10:00")
    add_subtitle(s, "Avant toute commande : comprendre le périmètre.")
    qs = [
        ("Qui ?", "1 utilisateur, un service, un étage, toute l’entreprise ?"),
        ("Quoi ?", "Pas Internet ? Lent ? Un seul site ? Mail OK, web KO ?"),
        ("Où ?", "Wi-Fi / filaire / VPN / site distant ?"),
        ("Quand ?", "Depuis 5 min, ce matin, après une MAJ / un orage ?"),
        ("Changement ?", "Nouveau switch, VLAN, firewall, box FAI, GPO ?"),
    ]
    y = 1.85
    for q, a in qs:
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(0.85))
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT
        c.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        tf = c.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = q + "  "
        set_run(run, 18, True, ACCENT)
        run = p.add_run()
        run.text = a
        set_run(run, 16, False, DARK)
        y += 0.95
    add_footer(s, 6)

    # 7 ETAPE 2
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Étape 2 — Couche 1 (physique)")
    add_timing(s, "10:00 – 16:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "À vérifier",
             "• Câble / connecteur RJ45\n"
             "• LED du port switch (up/down)\n"
             "• Bon SSID Wi-Fi (pas l’invité)\n"
             "• Airplane mode / Wi-Fi off\n"
             "• Port administratively down\n"
             "• Adaptateur désactivé (OS)\n"
             "• Alim AP / switch\n\n"
             "70 % des « pannes réseau »\ncommencent ici.",
             LIGHT, NAVY, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Commandes / gestes",
             "Windows\n"
             "  ncpa.cpl\n"
             "  (carte activée ?)\n\n"
             "Linux\n"
             "  ip link\n"
             "  nmcli device status\n\n"
             "Switch\n"
             "  show interface status\n"
             "  show interfaces fa0/1",
             GREEN_BG, GREEN, 14)
    add_footer(s, 7)

    # 8 ETAPE 3
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Étape 3 — Couche 2 (lien / VLAN / ARP)")
    add_timing(s, "16:00 – 22:00")
    add_bullets(s, [
        "Lien up mais mauvais VLAN → IP OK, GW KO ou mauvais réseau",
        "ARP : la machine connaît-elle la MAC de la passerelle ?",
        "AP isolé / client isolation Wi-Fi → ping voisin KO",
        "Boucle / STP : saturations, clignotements, réseau « mort »",
        "Mauvais trunk : VLAN manquant sur le trajet",
    ], top=1.65, size=17)
    add_code_card(
        s, 0.5, 4.5, 12.3, 2.1,
        "Commandes utiles",
        "arp -a                          # Windows / Linux\n"
        "show mac address-table           # Switch\n"
        "show vlan brief / show interfaces trunk",
    )
    add_footer(s, 8)

    # 9 ETAPE 4
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Étape 4 — Couche 3 (IP / masque / GW)")
    add_timing(s, "22:00 – 28:00")
    add_subtitle(s, "Rappel vidéo 1 : IP, masque et passerelle doivent être cohérents.")
    add_card(s, 0.5, 1.9, 4.0, 4.4, "Symptômes",
             "• 169.254.x.x (APIPA)\n"
             "• Pas de passerelle\n"
             "• GW hors subnet\n"
             "• Masque faux\n"
             "• Conflit d’IP\n"
             "• Double DHCP",
             RED_BG, RED, 14)
    add_card(s, 4.7, 1.9, 4.0, 4.4, "Actions",
             "• ipconfig /all  ou  ip a\n"
             "• Ping de la GW\n"
             "• Vérifier DHCP vs fixe\n"
             "• Renew bail DHCP\n"
             "• Table de routage\n"
             "  (route print / ip route)",
             LIGHT, NAVY, 14)
    add_card(s, 8.9, 1.9, 3.9, 4.4, "Règle",
             "GW dans le même\nsous-réseau que l’IP.\n\n"
             "Sinon : LAN parfois OK\nentre voisins, mais\njamais d’Internet.",
             YELLOW_BG, ORANGE, 14)
    add_footer(s, 9)

    # 10 ETAPE 5
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Étape 5 — DNS vs IP publique")
    add_timing(s, "28:00 – 34:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Le test décisif",
             "ping 8.8.8.8      → OK\n"
             "ping google.com  → KO\n\n"
             "= Problème DNS\n  (pas le câble, pas la GW)\n\n"
             "Inversement :\n"
             "ping 8.8.8.8 KO\n"
             "= Routage / NAT / FAI / FW\n  (le DNS n’est pas encore\n   le premier suspect)",
             LIGHT, NAVY, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Commandes DNS",
             "nslookup google.com\ndig google.com\n\n"
             "Windows\n"
             "  ipconfig /flushdns\n"
             "  (vider le cache)\n\n"
             "Tester un autre DNS\n"
             "  1.1.1.1  ou  8.8.8.8\n"
             "(temporairement, lab)",
             GREEN_BG, GREEN, 14)
    add_footer(s, 10)

    # 11 ETAPE 6
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Étape 6 — Ports, firewall, proxy")
    add_timing(s, "34:00 – 40:00")
    add_bullets(s, [
        "Ping OK ≠ application OK (ICMP ≠ HTTPS)",
        "Tester le port : 443, 80, 25, 3389, 22…",
        "Firewall local (Windows Defender) / firewall d’entreprise",
        "Proxy web obligatoire : le navigateur marche, curl non (ou l’inverse)",
        "Certificat TLS / date système fausse → sites HTTPS KO",
        "Filtrage DNS / category web / géobloc",
    ], top=1.65, size=17)
    add_code_card(
        s, 0.5, 4.6, 12.3, 2.0,
        "Tests ports",
        "Test-NetConnection google.com -Port 443     # PowerShell\n"
        "curl -I https://exemple.com\n"
        "nc -vz exemple.com 443                      # Linux",
    )
    add_footer(s, 11)

    # 12 ETAPE 7
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Étape 7 — Journaliser et escalader")
    add_timing(s, "40:00 – 44:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Ce que tu notes",
             "• Heure de début / fin\n"
             "• Symptôme exact (citation user)\n"
             "• Périmètre (1 PC / VLAN / site)\n"
             "• Commandes + résultats\n"
             "• Changements effectués\n"
             "• Hypothèse retenue\n"
             "• Capture éventuelle (Wireshark)\n\n"
             "Sans notes = panne qui revient\net personne ne sait pourquoi.",
             LIGHT, NAVY, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Quand escalader",
             "• FAI / box opérateur\n"
             "• Équipe sécu (firewall)\n"
             "• Cloud / SaaS (status page)\n"
             "• Fournisseur VPN\n\n"
             "Escalader AVEC un résumé\nclair des 7 étapes,\npas juste « ça marche pas ».",
             YELLOW_BG, ORANGE, 14)
    add_footer(s, 12)

    # 13 4 PINGS
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Les 4 ping (à graver)")
    add_timing(s, "Méthode")
    tests = [
        ("1", "127.0.0.1", "Pile TCP/IP locale", "Si KO → OS / stack"),
        ("2", "Passerelle", "LAN + config L3", "Si KO → câble, VLAN, IP, GW"),
        ("3", "8.8.8.8", "Sortie Internet", "Si KO → NAT / FAI / FW"),
        ("4", "google.com", "Résolution DNS", "Si KO et 3 OK → DNS"),
    ]
    y = 1.65
    for num, target, meaning, tip in tests:
        n = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(0.6), Inches(1.05))
        n.fill.solid()
        n.fill.fore_color.rgb = NAVY
        n.line.fill.background()
        n.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = n.text_frame.paragraphs[0].add_run()
        run.text = "\n" + num
        set_run(run, 20, True, WHITE)
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(y), Inches(11.4), Inches(1.05))
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT
        c.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"ping {target}   —   {meaning}"
        set_run(run, 16, True, NAVY)
        p2 = tf.add_paragraph()
        run = p2.add_run()
        run.text = tip
        set_run(run, 14, False, DARK)
        y += 1.2
    add_footer(s, 13)

    # 14 MATRICE
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Matrice rapide : lire les 4 ping")
    add_timing(s, "Aide-mémoire")
    rows = [
        ("1", "2", "3", "4", "Diagnostic probable"),
        ("KO", "—", "—", "—", "Stack OS / adaptateur"),
        ("OK", "KO", "—", "—", "Câble, VLAN, IP, GW"),
        ("OK", "OK", "KO", "—", "NAT, WAN, FAI, firewall"),
        ("OK", "OK", "OK", "KO", "DNS / filtre DNS"),
        ("OK", "OK", "OK", "OK", "Appliquer étape 6 (port/proxy/cert)"),
    ]
    y = 1.55
    widths = [1.1, 1.1, 1.1, 1.1, 7.5]
    xs = [0.5, 1.7, 2.9, 4.1, 5.3]
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            box = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(xs[j]), Inches(y), Inches(widths[j]), Inches(0.75)
            )
            if i == 0:
                box.fill.solid()
                box.fill.fore_color.rgb = NAVY
                col = WHITE
                bold = True
            else:
                box.fill.solid()
                box.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
                col = DARK
                bold = j == 4
            box.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
            box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if j < 4 else PP_ALIGN.LEFT
            run = box.text_frame.paragraphs[0].add_run()
            run.text = (("  " + val) if j == 4 else val)
            set_run(run, 13, bold, col)
        y += 0.75
    add_footer(s, 14)

    # 15 INCIDENT A
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Incident A — Pas de lien")
    add_timing(s, "44:00 – 48:00")
    add_subtitle(s, "Symptôme : icône réseau débranché / « Network cable unplugged »")
    add_card(s, 0.5, 1.9, 4.0, 4.4, "Cadrage",
             "• 1 seul PC\n"
             "• Filaire\n"
             "• Depuis ce matin\n"
             "• Hier OK\n"
             "• Agent a « rangé »\n"
             "  les câbles hier soir",
             LIGHT, NAVY, 14)
    add_card(s, 4.7, 1.9, 4.0, 4.4, "Parcours",
             "1. Cadrer → 1 PC\n"
             "2. Physique → LED off\n"
             "3. Tester autre câble\n"
             "4. Autre port switch\n"
             "5. show int status\n"
             "6. Port shutdown\n"
             "   trouvé",
             YELLOW_BG, ORANGE, 14)
    add_card(s, 8.9, 1.9, 3.9, 4.4, "Résolution",
             "no shutdown\nsur le port\n\n"
             "Documenter :\nqui a shutdown\net pourquoi\n\n"
             "Leçon : couche 1\navant ping google",
             GREEN_BG, GREEN, 14)
    add_footer(s, 15)

    # 16 INCIDENT B
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Incident B — LAN OK, Internet KO")
    add_timing(s, "48:00 – 53:00")
    add_subtitle(s, "Symptôme : partage de fichiers OK, web et mail externes KO")
    add_code_card(
        s, 0.5, 1.9, 6.2, 4.5,
        "Résultats des 4 ping",
        "ping 127.0.0.1      OK\n"
        "ping 192.168.1.1    OK   (GW)\n"
        "ping 8.8.8.8        KO\n"
        "ping google.com     KO\n\n"
        "→ Étape 5 inutile tant que 3 = KO\n"
        "→ Suspect : NAT / WAN / FAI / FW",
    )
    add_card(s, 7.0, 1.9, 5.8, 4.5, "Actions",
             "• Traceroute / tracert 8.8.8.8\n"
             "• Vérifier WAN box / routeur\n"
             "• Lien FAI up ?\n"
             "• Règle firewall outbound\n"
             "• Autre réseau (tél. 4G)\n"
             "  pour isoler le site\n\n"
             "Cause fréquente lab :\n"
             "NAT oublié / WAN down",
             LIGHT, NAVY, 14)
    add_footer(s, 16)

    # 17 INCIDENT C
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Incident C — Un site marche, l’autre non")
    add_timing(s, "53:00 – 58:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Indices",
             "• 4 ping tous OK\n"
             "• google.com OK\n"
             "• intranet.entreprise.local KO\n"
             "  ou un SaaS précis KO\n\n"
             "Donc : pas le câble,\npas la GW, pas « Internet ».\n\n"
             "→ Étape 6 + DNS interne",
             LIGHT, NAVY, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Pistes",
             "• DNS interne vs public\n"
             "• Proxy / authentification\n"
             "• Certificat expiré\n"
             "• Filtrage catégorie web\n"
             "• Firewall port 443 vers\n"
             "  cette destination\n"
             "• Service distant down\n"
             "  (status page)\n\n"
             "nslookup + curl -I + navigateur",
             YELLOW_BG, ORANGE, 14)
    add_footer(s, 17)

    # 18 FICHE RECAP
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Fiche récap (screenshot)")
    add_timing(s, "58:00")
    add_card(s, 0.5, 1.6, 6.0, 5.0, "Ordre d’attaque",
             "1. Cadrer (qui/quoi/où/quand)\n"
             "2. Physique (câble, LED, Wi-Fi)\n"
             "3. Lien (VLAN, ARP, trunk)\n"
             "4. IP / masque / GW\n"
             "5. 8.8.8.8 vs nom DNS\n"
             "6. Port / FW / proxy / cert\n"
             "7. Notes + escalade\n\n"
             "Une couche à la fois.",
             LIGHT, NAVY, 14)
    add_card(s, 6.8, 1.6, 5.9, 5.0, "4 ping",
             "1. 127.0.0.1\n"
             "2. Passerelle\n"
             "3. 8.8.8.8 (ou IP WAN)\n"
             "4. google.com\n\n"
             "Windows\n"
             "  ipconfig /all\n"
             "  ping …\n"
             "  nslookup …\n"
             "  Test-NetConnection -Port 443",
             GREEN_BG, GREEN, 14)
    add_footer(s, 18)

    # 19 OSI RAPPEL
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "OSI en mode dépannage (pas théorique)")
    add_timing(s, "Astuce")
    layers = [
        ("7–5", "Application", "Navigateur, mail, cert, proxy"),
        ("4", "Transport", "Ports TCP/UDP, firewall L4"),
        ("3", "Réseau", "IP, masque, GW, routage"),
        ("2", "Liaison", "MAC, VLAN, switch, ARP"),
        ("1", "Physique", "Câble, Wi-Fi radio, LED"),
    ]
    y = 1.6
    for num, name, use in layers:
        left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(1.8), Inches(0.9))
        left.fill.solid()
        left.fill.fore_color.rgb = NAVY
        left.line.fill.background()
        left.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = left.text_frame.paragraphs[0].add_run()
        run.text = "\n" + num
        set_run(run, 16, True, WHITE)
        mid = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(y), Inches(3.5), Inches(0.9))
        mid.fill.solid()
        mid.fill.fore_color.rgb = ACCENT
        mid.line.fill.background()
        mid.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = mid.text_frame.paragraphs[0].add_run()
        run.text = "\n" + name
        set_run(run, 16, True, WHITE)
        right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), Inches(y), Inches(6.6), Inches(0.9))
        right.fill.solid()
        right.fill.fore_color.rgb = LIGHT
        right.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        run = right.text_frame.paragraphs[0].add_run()
        run.text = "  " + use
        set_run(run, 16, False, DARK)
        y += 1.0
    add_footer(s, 19)

    # 20 ERREURS
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Erreurs de diagnostic (à éviter)")
    add_timing(s, "Pièges")
    errs = [
        "Changer DNS + câble + firewall en même temps",
        "Conclure « c’est le DNS » sans avoir testé 8.8.8.8",
        "Ignorer qu’un seul user est touché (problème local)",
        "Oublier le VPN / le proxy dans le cadrage",
        "Ne rien écrire → la panne revient la semaine suivante",
        "Blâmer le FAI avant d’avoir ping la GW",
    ]
    y = 1.6
    for e in errs:
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(0.75))
        c.fill.solid()
        c.fill.fore_color.rgb = RED_BG
        c.line.color.rgb = RED
        run = c.text_frame.paragraphs[0].add_run()
        run.text = "✖  " + e
        set_run(run, 16, False, DARK)
        y += 0.85
    add_footer(s, 20)

    # 21 RECAP
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Récap — ce que tu repars avec")
    add_timing(s, "59:00")
    recaps = [
        ("Méthode", "7 étapes, du cadrage à l’escalade"),
        ("4 ping", "Localhost → GW → IP publique → DNS"),
        ("OSI utile", "Physique → … → application, une couche à la fois"),
        ("3 incidents", "Lien / Internet / un seul site"),
        ("Pro", "Notes + hypothèse + une modif à la fois"),
    ]
    y = 1.65
    for a, b in recaps:
        left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(3.3), Inches(0.88))
        left.fill.solid()
        left.fill.fore_color.rgb = NAVY
        left.line.fill.background()
        left.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = left.text_frame.paragraphs[0].add_run()
        run.text = "\n" + a
        set_run(run, 16, True, WHITE)
        right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(y), Inches(8.8), Inches(0.88))
        right.fill.solid()
        right.fill.fore_color.rgb = LIGHT
        right.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        run = right.text_frame.paragraphs[0].add_run()
        run.text = "  " + b
        set_run(run, 16, False, DARK)
        y += 1.0
    add_footer(s, 21)

    # 22 DEVOIR + CTA
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Devoir + fin de playlist")
    add_timing(s, "59:20 – 60:00")
    add_card(s, 0.5, 1.7, 6.0, 4.6, "Devoir (commentaire)",
             "Décris UNE panne réelle\n(ou inventée) avec :\n\n"
             "• Cadrage (qui/quoi/où/quand)\n"
             "• Résultat des 4 ping\n"
             "• Étape où tu as trouvé\n"
             "• Solution\n\n"
             "Je corrige les meilleures\nen commentaire épinglé.",
             YELLOW_BG, ORANGE, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.6, "Playlist « Réseau 1 h »",
             "V1 — IP, masque, passerelle\n"
             "V2 — VLAN + lab\n"
             "V3 — Diagnostic 7 étapes\n\n"
             "Tu as maintenant un socle\ndébutant → intermédiaire.\n\n"
             "Abonne-toi pour la suite\n(Wireshark / Firewall / VPN).",
             LIGHT, NAVY, 14)
    add_footer(s, 22)

    # 23 MERCI
    s = blank(prs)
    add_bg(s, NAVY)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Merci !"
    set_run(run, 48, True, WHITE)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(11.5), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Cadrer  •  4 ping  •  Une couche à la fois"
    set_run(run, 20, False, LIGHT)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "ABOU SAYABOU"
    set_run(run, 24, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run = p2.add_run()
    run.text = "Cours réseau — Vidéo 3 / 3 — Fin de la série"
    set_run(run, 14, False, ACCENT)

    # 24 NOTES
    s = blank(prs)
    add_top_bar(s, "Notes formateur (ne pas montrer)")
    add_title(s, "Notes de tournage")
    add_bullets(s, [
        "Garder la slide « 7 étapes » en favori pour y revenir pendant les incidents",
        "Démo réelle : débrancher un câble (incident A) si tu tournes en live",
        "Pour incident B : couper temporairement le NAT / WAN dans PT ou box",
        "Insister : ping OK ≠ HTTPS OK (étape 6)",
        "Rappeler V1 (GW) et V2 (VLAN) quand tu parles couche 2/3",
        "Devoir : épingler 1–2 meilleurs commentaires le lendemain",
        "Teaser prochaine série : Wireshark 1h / Firewall PME 1h",
        "Durée parlée cible : 55–60 min — ne pas griller le cadrage",
    ], top=1.65, size=16)
    add_footer(s, 24)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    stamp_meta(OUT)
    print(f"PowerPoint généré : {OUT} ({OUT.stat().st_size} octets) — {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
