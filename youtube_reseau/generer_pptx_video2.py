#!/usr/bin/env python3
"""PowerPoint — Vidéo 2 : VLAN + lab (1 h) — ABOU SAYABOU."""

from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Video2_VLAN_Lab_ABOU_SAYABOU.pptx"

NAVY = RGBColor(0x0D, 0x3B, 0x66)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
LIGHT = RGBColor(0xE8, 0xF1, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
RED = RGBColor(0xC6, 0x28, 0x28)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
VLAN10 = RGBColor(0x15, 0x65, 0xC0)  # RH
VLAN20 = RGBColor(0x2E, 0x7D, 0x32)  # Invités
VLAN30 = RGBColor(0xC6, 0x28, 0x28)  # Serveurs

W, H = Inches(13.333), Inches(7.5)
TOTAL = 24
FOOTER = "Vidéo 2 — VLAN + lab  |  Débutant → Intermédiaire"


def set_run(run, size=20, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_top_bar(slide, title_right="Cours réseau — 1 h"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    p = bar.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "  ABOU SAYABOU  •  Réseau informatique"
    set_run(run, 12, True, WHITE)
    lab = slide.shapes.add_textbox(Inches(8.5), Inches(0.12), Inches(4.5), Inches(0.35))
    p = lab.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = title_right + "  "
    set_run(run, 12, False, LIGHT)


def add_footer(slide, page):
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), W, Inches(0.35))
    foot.fill.solid()
    foot.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    foot.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.18), Inches(10), Inches(0.3))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = FOOTER
    set_run(run, 10, False, GRAY)
    tb2 = slide.shapes.add_textbox(Inches(11.5), Inches(7.18), Inches(1.5), Inches(0.3))
    p = tb2.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page}/{TOTAL}"
    set_run(run, 10, True, GRAY)


def add_title(slide, text, top=0.7):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.7))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    set_run(run, 30, True, NAVY)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top + 0.62), Inches(1.8), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_subtitle(slide, text, top=1.42):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.4))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    set_run(run, 16, False, GRAY)


def add_bullets(slide, items, left=0.6, top=2.0, width=12, size=18):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "▸  " + item
        set_run(run, size, False, DARK)


def add_card(slide, left, top, width, height, title, body, fill=LIGHT, title_color=NAVY, size=13):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    run = tf.paragraphs[0].add_run()
    run.text = title
    set_run(run, 16, True, title_color)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    run = p2.add_run()
    run.text = body
    set_run(run, size, False, DARK)
    return card


def add_timing(slide, text):
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.68), Inches(2.4), Inches(0.4)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    set_run(run, 12, True, WHITE)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_code_card(slide, left, top, width, height, title, code):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x32)
    card.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.1)
    run = tf.paragraphs[0].add_run()
    run.text = title
    set_run(run, 13, True, ACCENT, "Calibri")
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    run = p2.add_run()
    run.text = code
    set_run(run, 12, False, WHITE, "Consolas")


def stamp_meta(path: Path):
    now = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    buf = BytesIO(path.read_bytes())
    out_buf = BytesIO()
    with zipfile.ZipFile(buf, "r") as zin, zipfile.ZipFile(
        out_buf, "w", compression=zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "docProps/core.xml":
                data = f"""<?xml version='1.0' encoding='UTF-8' standalone='yes'?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <dc:title>Vidéo 2 — VLAN expliqué simplement + lab (cours 1h)</dc:title>
  <dc:subject>Cours réseau — VLAN, access, trunk, inter-VLAN</dc:subject>
  <dc:creator>ABOU SAYABOU</dc:creator>
  <cp:keywords>VLAN, 802.1Q, trunk, Packet Tracer, YouTube</cp:keywords>
  <dc:description>Support de cours vidéo 1h — VLAN + lab</dc:description>
  <cp:lastModifiedBy>ABOU SAYABOU</cp:lastModifiedBy>
  <cp:revision>1</cp:revision>
  <dcterms:created xsi:type="dcterms:W3CDTF">{now}</dcterms:created>
  <dcterms:modified xsi:type="dcterms:W3CDTF">{now}</dcterms:modified>
  <cp:category>Formation réseau</cp:category>
  <dc:language>fr-FR</dc:language>
</cp:coreProperties>""".encode("utf-8")
            elif item.filename == "docProps/app.xml":
                text = data.decode("utf-8")
                text = text.replace("<Company/>", "<Company>ABOU SAYABOU</Company>")
                text = text.replace("python-pptx", "ABOU SAYABOU")
                data = text.encode("utf-8")
            zout.writestr(item, data)
    path.write_bytes(out_buf.getvalue())


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 COVER
    s = blank(prs)
    add_bg(s, NAVY)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.8), W, Inches(1.7))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(0x0A, 0x2F, 0x52)
    strip.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "COURS RÉSEAU  •  VIDÉO 2 / 3"
    set_run(run, 18, True, ACCENT)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.2))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "VLAN expliqué simplement"
    set_run(run, 38, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    run = p2.add_run()
    run.text = "+ lab complet (3 VLAN, 1 routeur)"
    set_run(run, 28, True, LIGHT)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.5))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Durée 1 heure  •  Packet Tracer  •  Access • Trunk • Inter-VLAN"
    set_run(run, 16, False, RGBColor(0xB0, 0xC4, 0xDE))
    tb = s.shapes.add_textbox(Inches(0.8), Inches(6.15), Inches(11.5), Inches(0.8))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "ABOU SAYABOU"
    set_run(run, 22, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    run = p2.add_run()
    run.text = "Formateur cybersécurité & réseaux"
    set_run(run, 14, False, LIGHT)

    # 2 OBJECTIFS
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Objectif de la vidéo")
    add_timing(s, "0:00 – 3:00")
    add_subtitle(s, "À la fin de cette heure, tu sauras…")
    add_card(s, 0.5, 2.0, 3.9, 4.2, "1. Pourquoi VLAN",
             "Comprendre le problème du LAN plat : broadcasts, sécurité, Wi-Fi invités mélangés aux serveurs.",
             LIGHT)
    add_card(s, 4.7, 2.0, 3.9, 4.2, "2. Access vs Trunk",
             "Configurer des ports access (1 VLAN) et un trunk 802.1Q vers le routeur / l’autre switch.",
             LIGHT)
    add_card(s, 8.9, 2.0, 3.9, 4.2, "3. Lab qui marche",
             "3 VLAN (RH / Invités / Serveurs), inter-VLAN, tests ping, et une ACL pour isoler les invités.",
             LIGHT)
    add_footer(s, 2)

    # 3 PLAN
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Plan de la session (1 h)")
    add_timing(s, "Chapitres")
    plan = [
        ("01", "0:00", "Pourquoi les VLAN (le problème)"),
        ("02", "3:00", "Broadcast vs collision — rappel"),
        ("03", "8:00", "VLAN = LAN logique"),
        ("04", "13:00", "Access vs Trunk (802.1Q)"),
        ("05", "20:00", "Native VLAN et erreurs classiques"),
        ("06", "26:00", "Inter-VLAN : SVI vs router-on-a-stick"),
        ("07", "34:00", "Lab Packet Tracer — topologie"),
        ("08", "40:00", "Config switch + trunk + routeur + tests"),
        ("09", "57:30", "Sécu invités, récap, devoir"),
    ]
    y = 1.5
    for num, t, label in plan:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(0.7), Inches(0.48))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = box.text_frame.paragraphs[0].add_run()
        run.text = num
        set_run(run, 14, True, WHITE)
        tbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.35), Inches(y), Inches(1.3), Inches(0.48))
        tbox.fill.solid()
        tbox.fill.fore_color.rgb = ACCENT
        tbox.line.fill.background()
        tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tbox.text_frame.paragraphs[0].add_run()
        run.text = t
        set_run(run, 13, True, WHITE)
        tb = s.shapes.add_textbox(Inches(2.9), Inches(y + 0.05), Inches(9.5), Inches(0.4))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = label
        set_run(run, 18, False, DARK)
        y += 0.55
    add_footer(s, 3)

    # 4 POURQUOI
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Le problème : un LAN plat")
    add_timing(s, "0:00 – 3:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Sans VLAN",
             "• 1 switch = 1 grand réseau\n"
             "• Tout le monde entend les broadcasts\n"
             "• RH, invités Wi-Fi, serveurs mélangés\n"
             "• Un PC infecté scanne tout le LAN\n"
             "• Imprimantes et caméras exposées\n"
             "• Difficile à faire grandir (scaling)\n\n"
             "C’est le réseau « ça marche »…\njusqu’au premier incident.",
             RGBColor(0xFF, 0xEB, 0xEE), RED)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Avec VLAN",
             "• Même switch physique\n"
             "• Plusieurs LAN logiques\n"
             "• Broadcasts contenus par VLAN\n"
             "• Isolation : invités ≠ serveurs\n"
             "• Politiques firewall plus simples\n"
             "• Base d’un LAN PME propre\n\n"
             "Objectif du lab : 3 mondes\nsur 1 seul switch.",
             RGBColor(0xE8, 0xF5, 0xE9), GREEN)
    add_footer(s, 4)

    # 5 BROADCAST
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Rappel : collision vs broadcast")
    add_timing(s, "3:00 – 8:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Domaine de collision",
             "• Qui « parle » en même temps sur un média\n"
             "• Switch moderne : 1 collision par port\n"
             "  (full duplex)\n"
             "• Moins un sujet aujourd’hui\n\n"
             "À retenir : le switch sépare\nles collisions.",
             LIGHT, NAVY)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Domaine de broadcast",
             "• Qui reçoit un paquet « à tout le monde »\n"
             "  (ARP, DHCP discover, etc.)\n"
             "• Par défaut : TOUT le VLAN / le LAN\n"
             "• Le routeur / SVI arrête le broadcast\n\n"
             "VLAN = découper le domaine\nde broadcast.",
             RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
    add_footer(s, 5)

    # 6 VLAN DEFINITION
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "VLAN = LAN logique")
    add_timing(s, "8:00 – 13:00")
    add_subtitle(s, "Même matériel, 3 réseaux. IEEE 802.1Q.")
    add_card(s, 0.5, 1.9, 3.9, 4.4, "VLAN 10 — RH",
             "10.10.10.0/24\nGW : 10.10.10.1\n\nPC RH, imprimantes\ninterne, téléphonie\nde bureau.",
             RGBColor(0xE3, 0xF2, 0xFD), VLAN10)
    add_card(s, 4.7, 1.9, 3.9, 4.4, "VLAN 20 — Invités",
             "10.10.20.0/24\nGW : 10.10.20.1\n\nWi-Fi guests,\nvisiteurs.\nPas d’accès serveurs.",
             RGBColor(0xE8, 0xF5, 0xE9), VLAN20)
    add_card(s, 8.9, 1.9, 3.9, 4.4, "VLAN 30 — Serveurs",
             "10.10.30.0/24\nGW : 10.10.30.1\n\nAD, fichiers, NAS,\napplications métiers.",
             RGBColor(0xFF, 0xEB, 0xEE), VLAN30)
    add_footer(s, 6)

    # 7 ACCESS TRUNK
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Access vs Trunk")
    add_timing(s, "13:00 – 20:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Port ACCESS",
             "• 1 seul VLAN (untagged)\n"
             "• Branché à un PC, imprimante, AP\n"
             "• Le PC ne « voit » pas le tag 802.1Q\n\n"
             "switchport mode access\n"
             "switchport access vlan 10\n\n"
             "Erreur : mettre un PC en trunk\n→ parfois ça marche… mal.",
             LIGHT, NAVY)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Port TRUNK",
             "• Plusieurs VLAN (tagged 802.1Q)\n"
             "• Switch ↔ switch  ou  switch ↔ routeur\n"
             "• Chaque trame porte un VLAN ID\n\n"
             "switchport mode trunk\n"
             "switchport trunk allowed vlan 10,20,30\n\n"
             "Toujours limiter les VLAN\nautorisés (pas « all » en prod).",
             RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
    add_footer(s, 7)

    # 8 802.1Q visuel
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Le tag 802.1Q (idée simple)")
    add_timing(s, "Schéma")
    add_bullets(s, [
        "Sur un lien trunk, le switch ajoute 4 octets : le VLAN ID (1–4094)",
        "VLAN 10, 20, 30 voyagent sur LE MÊME câble vers le routeur",
        "À l’arrivée, le routeur « décapsule » et route d’un VLAN vers un autre",
        "Sur un port access : pas de tag — le VLAN est implicite",
    ], top=1.7, size=18)
    # fake frame
    labels = [("Dest MAC", 2.0), ("Src MAC", 2.0), ("802.1Q\nVLAN 10", 2.4), ("Type", 1.5), ("Données…", 3.5)]
    x = 0.7
    colors = [NAVY, NAVY, ACCENT, GRAY, GREEN]
    y = 4.3
    for (lab, w), col in zip(labels, colors):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = col
        box.line.fill.background()
        box.text_frame.word_wrap = True
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = box.text_frame.paragraphs[0].add_run()
        run.text = lab
        set_run(run, 13, True, WHITE)
        x += w + 0.12
    add_footer(s, 8)

    # 9 NATIVE VLAN
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Native VLAN et erreurs classiques")
    add_timing(s, "20:00 – 26:00")
    add_bullets(s, [
        "Native VLAN = VLAN des trames SANS tag sur le trunk (souvent VLAN 1 par défaut)",
        "En production : ne pas laisser le VLAN 1 pour la data ni pour le native",
        "Les deux extrémités du trunk doivent avoir le MÊME native VLAN",
        "Mismatch access / trunk = lien « bizarre », CDP native VLAN mismatch",
        "Oublier allowed vlan → trop de VLAN transitent (attaque VLAN hopping / bruit)",
    ], top=1.7, size=18)
    tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.2))
    tip.fill.solid()
    tip.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    tip.line.color.rgb = RED
    tf = tip.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = "Règle d’or : un PC = port ACCESS. Un lien switch/routeur = TRUNK. Jamais l’inverse « par habitude »."
    set_run(run, 16, True, RED)
    add_footer(s, 9)

    # 10 INTERVLAN
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Inter-VLAN routing : 2 méthodes")
    add_timing(s, "26:00 – 34:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Router-on-a-stick",
             "• 1 lien physique trunk vers un routeur\n"
             "• Sous-interfaces : G0/0.10 .20 .30\n"
             "• encapsulation dot1Q 10 / 20 / 30\n"
             "• Simple pour un lab / petite PME\n"
             "• Le lien unique peut saturer\n\n"
             "→ C’est CE qu’on lab aujourd’hui",
             LIGHT, NAVY)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "SVI (switch L3)",
             "• Switch multilayer (ex. 3560, 9300)\n"
             "• interface vlan 10 + ip address\n"
             "• ip routing\n"
             "• Plus propre, plus performant\n"
             "• Souvent + ACL / firewall ensuite\n\n"
             "→ À mentionner : « la vraie vie PME »",
             RGBColor(0xE8, 0xF5, 0xE9), GREEN)
    add_footer(s, 10)

    # 11 TOPO LAB
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Lab : topologie Packet Tracer")
    add_timing(s, "34:00 – 40:00")
    add_subtitle(s, "3 PC  •  1 switch L2  •  1 routeur  •  1 câble trunk")

    nodes = [
        (0.4, "PC-RH\n.10\nVLAN 10", VLAN10),
        (3.0, "PC-INV\n.10\nVLAN 20", VLAN20),
        (5.6, "PC-SRV\n.10\nVLAN 30", VLAN30),
        (8.2, "SW1\nFa0/24 trunk", GRAY),
        (10.8, "R1\nGW .1", ORANGE),
    ]
    for x, text, col in nodes:
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.15), Inches(2.15), Inches(1.7))
        sh.fill.solid()
        sh.fill.fore_color.rgb = col
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = text
        set_run(run, 13, True, WHITE)

    add_card(s, 0.5, 4.15, 4.0, 2.3, "Adressage",
             "VLAN10  10.10.10.0/24\nVLAN20  10.10.20.0/24\nVLAN30  10.10.30.0/24",
             LIGHT, size=14)
    add_card(s, 4.7, 4.15, 4.0, 2.3, "Ports switch",
             "Fa0/1  access VLAN 10\nFa0/2  access VLAN 20\nFa0/3  access VLAN 30\nFa0/24 trunk → R1",
             LIGHT, size=14)
    add_card(s, 8.9, 4.15, 3.9, 2.3, "PC",
             "IP .10 de leur VLAN\nMasque /24\nGW = .1 du VLAN",
             LIGHT, size=14)
    add_footer(s, 11)

    # 12 CONFIG VLAN SWITCH
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Config switch — VLAN et access")
    add_timing(s, "40:00 – 48:00")
    add_code_card(
        s, 0.5, 1.55, 12.3, 5.0,
        "SW1",
        "enable\n"
        "configure terminal\n"
        "vlan 10\n name RH\n"
        "vlan 20\n name INVITES\n"
        "vlan 30\n name SERVEURS\n"
        "interface range fa0/1\n switchport mode access\n switchport access vlan 10\n switchport nonegotiate\n"
        "interface fa0/2\n switchport mode access\n switchport access vlan 20\n"
        "interface fa0/3\n switchport mode access\n switchport access vlan 30",
    )
    add_footer(s, 12)

    # 13 TRUNK
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Config switch — trunk")
    add_timing(s, "Trunk")
    add_code_card(
        s, 0.5, 1.55, 7.2, 5.0,
        "SW1 — Fa0/24 vers le routeur",
        "interface fa0/24\n"
        " switchport trunk encapsulation dot1q\n"
        " switchport mode trunk\n"
        " switchport trunk allowed vlan 10,20,30\n"
        " switchport trunk native vlan 99\n"
        " no shutdown\n\n"
        "do show vlan brief\n"
        "do show interfaces trunk",
    )
    add_card(s, 8.0, 1.55, 4.8, 5.0, "Vérifs",
             "show vlan brief\n→ ports dans le bon VLAN\n\n"
             "show interfaces trunk\n→ Fa0/24 trunking\n\n"
             "Native VLAN 99 =\nVLAN dummy (pas de data)\n\n"
             "Si encapsulation n’existe pas\n(switch L2 PT) :\nomets la ligne encapsulation.",
             LIGHT, size=14)
    add_footer(s, 13)

    # 14 ROUTEUR
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Config routeur — router-on-a-stick")
    add_timing(s, "48:00 – 54:00")
    add_code_card(
        s, 0.5, 1.55, 12.3, 5.0,
        "R1 — sous-interfaces",
        "configure terminal\n"
        "interface g0/0\n no shutdown\n no ip address\n"
        "interface g0/0.10\n encapsulation dot1Q 10\n ip address 10.10.10.1 255.255.255.0\n"
        "interface g0/0.20\n encapsulation dot1Q 20\n ip address 10.10.20.1 255.255.255.0\n"
        "interface g0/0.30\n encapsulation dot1Q 30\n ip address 10.10.30.1 255.255.255.0\n"
        "exit\n"
        "do show ip route\n"
        "do show ip interface brief",
    )
    add_footer(s, 14)

    # 15 PC
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Config des PC (IP fixe pour le lab)")
    add_timing(s, "PC")
    add_card(s, 0.5, 1.7, 3.9, 4.6, "PC-RH",
             "IP      10.10.10.10\nMasque  255.255.255.0\nGW      10.10.10.1\nDNS     10.10.10.1\n\nPort SW : Fa0/1\nVLAN 10",
             RGBColor(0xE3, 0xF2, 0xFD), VLAN10, 14)
    add_card(s, 4.7, 1.7, 3.9, 4.6, "PC-INV",
             "IP      10.10.20.10\nMasque  255.255.255.0\nGW      10.10.20.1\nDNS     10.10.20.1\n\nPort SW : Fa0/2\nVLAN 20",
             RGBColor(0xE8, 0xF5, 0xE9), VLAN20, 14)
    add_card(s, 8.9, 1.7, 3.9, 4.6, "PC-SRV",
             "IP      10.10.30.10\nMasque  255.255.255.0\nGW      10.10.30.1\nDNS     10.10.30.1\n\nPort SW : Fa0/3\nVLAN 30",
             RGBColor(0xFF, 0xEB, 0xEE), VLAN30, 14)
    add_footer(s, 15)

    # 16 TESTS
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Tests ping — qui passe ?")
    add_timing(s, "54:00 – 57:30")
    rows = [
        ("PC-RH → 10.10.10.1 (GW)", "OK", "Même VLAN + GW"),
        ("PC-RH → PC-INV (10.10.20.10)", "OK", "Inter-VLAN via R1"),
        ("PC-RH → PC-SRV (10.10.30.10)", "OK", "Inter-VLAN (à restreindre)"),
        ("PC-RH → 10.10.20.10 sans trunk", "KO", "Pas de chemin L3"),
        ("PC access VLAN 10 ping .20.10\nsans IP sur .20 du routeur", "KO", "Pas de GW inter-VLAN"),
    ]
    y = 1.55
    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(0.42))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = NAVY
    hdr.line.fill.background()
    run = hdr.text_frame.paragraphs[0].add_run()
    run.text = "    Test                                      Résultat     Pourquoi"
    set_run(run, 14, True, WHITE)
    y = 2.05
    for i, (test, res, why) in enumerate(rows):
        bg = LIGHT if i % 2 == 0 else WHITE
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(0.78))
        box.fill.solid()
        box.fill.fore_color.rgb = bg
        box.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = test + "   →   "
        set_run(run, 14, False, DARK)
        run = p.add_run()
        run.text = res
        set_run(run, 14, True, GREEN if res == "OK" else RED)
        run = p.add_run()
        run.text = "   •  " + why
        set_run(run, 14, False, GRAY)
        y += 0.82
    add_footer(s, 16)

    # 17 ACL
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Isolation invités : ACL (bonus lab)")
    add_timing(s, "Sécu")
    add_subtitle(s, "Les VLAN isolent L2. Pour bloquer Invités → Serveurs, il faut du L3 (ACL / firewall).")
    add_code_card(
        s, 0.5, 1.95, 7.4, 4.5,
        "R1 — ACL",
        "ip access-list extended BLOQUE-INVITES\n"
        " deny ip 10.10.20.0 0.0.0.255 10.10.30.0 0.0.0.255\n"
        " permit ip any any\n"
        "interface g0/0.20\n"
        " ip access-group BLOQUE-INVITES in\n\n"
        "! Test : PC-INV ping PC-SRV  →  KO\n"
        "! Test : PC-INV ping GW     →  OK",
    )
    add_card(s, 8.15, 1.95, 4.65, 4.5, "Message clé",
             "VLAN ≠ firewall.\n\n"
             "VLAN = séparation\ndes broadcasts.\n\n"
             "ACL / firewall =\nqui a le droit d’aller\noù.\n\n"
             "En prod : plutôt\nfirewall (FortiGate)\nentre les zones.",
             RGBColor(0xFF, 0xF3, 0xE0), ORANGE, 14)
    add_footer(s, 17)

    # 18 SECU
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Sécurité VLAN — checklist PME")
    add_timing(s, "57:30 – 59:00")
    checks = [
        "Pas de data sur VLAN 1 (change native + management dédié)",
        "Ports access : nonegotiate (coupe DTP)",
        "Trunk : allowed vlan list explicite",
        "Ports inutilisés : shutdown + VLAN parking (ex. 999)",
        "Invités : pas d’accès VLAN serveurs (ACL ou firewall)",
        "Wi-Fi invité = VLAN 20 + isolation client AP si possible",
        "Documenter : VLAN ID, nom, subnet, GW, usage",
    ]
    add_bullets(s, checks, top=1.65, size=17)
    add_footer(s, 18)

    # 19 ERREURS LAB
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Si le ping ne passe pas…")
    add_timing(s, "Dépannage lab")
    errs = [
        ("PC et GW pas le même VLAN / subnet", "Revérifier IP .10 et GW .1"),
        ("Port encore en VLAN 1", "show vlan brief"),
        ("Trunk down / pas de dot1Q", "show interfaces trunk"),
        ("Sous-interface routeur shutdown", "show ip interface brief"),
        ("Câble sur le mauvais port", "Topo Packet Tracer"),
        ("ACL trop large (deny any)", "Retirer l’ACL pour tester"),
    ]
    y = 1.6
    for title, cmd in errs:
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(0.78))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
        c.line.color.rgb = ORANGE
        tf = c.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "⚠  " + title + "   →   "
        set_run(run, 15, True, ORANGE)
        run = p.add_run()
        run.text = cmd
        set_run(run, 15, False, DARK)
        y += 0.85
    add_footer(s, 19)

    # 20 RECAP
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Récap — 5 phrases")
    add_timing(s, "59:00")
    recaps = [
        ("VLAN", "Découpe le domaine de broadcast (LAN logique)"),
        ("Access", "Un port, un VLAN — pour les terminaux"),
        ("Trunk", "Plusieurs VLAN taggés 802.1Q — pour switch/routeur"),
        ("Inter-VLAN", "Il faut un routeur (stick) ou un SVI L3"),
        ("Sécu", "VLAN ≠ firewall : ACL / FW pour invités → serveurs"),
    ]
    y = 1.65
    for a, b in recaps:
        left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(3.3), Inches(0.88))
        left.fill.solid()
        left.fill.fore_color.rgb = NAVY
        left.line.fill.background()
        left.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = left.text_frame.paragraphs[0].add_run()
        run.text = "\n" + a
        set_run(run, 16, True, WHITE)
        right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(y), Inches(8.8), Inches(0.88))
        right.fill.solid()
        right.fill.fore_color.rgb = LIGHT
        right.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        run = right.text_frame.paragraphs[0].add_run()
        run.text = "  " + b
        set_run(run, 16, False, DARK)
        y += 1.0
    add_footer(s, 20)

    # 21 HOMEWORK + CTA
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Devoir + prochaine vidéo")
    add_timing(s, "59:00 – 60:00")
    add_card(s, 0.5, 1.7, 6.0, 4.6, "Devoir (commentaire)",
             "Ajoute un VLAN 40 VoIP :\n\n"
             "• subnet 10.10.40.0/24\n"
             "• GW 10.10.40.1\n"
             "• un PC-VOIP .10\n"
             "• commande bonus :\n"
             "  switchport voice vlan 40\n\n"
             "Poste un screenshot show vlan brief !",
             RGBColor(0xFF, 0xF3, 0xE0), ORANGE, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.6, "Vidéo 3",
             "Le réseau est down :\n7 étapes pour diagnostiquer\nn’importe quelle panne\n\n"
             "Méthode OSI + 3 incidents\n+ les 4 ping de la vidéo 1\n\n"
             "Abonne-toi à la playlist\n« Réseau 1 h »",
             LIGHT, NAVY, 14)
    add_footer(s, 21)

    # 22 MERCI
    s = blank(prs)
    add_bg(s, NAVY)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Merci !"
    set_run(run, 48, True, WHITE)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11.5), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "VLAN  •  Access  •  Trunk  •  Inter-VLAN  •  ACL"
    set_run(run, 20, False, LIGHT)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "ABOU SAYABOU"
    set_run(run, 24, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run = p2.add_run()
    run.text = "Cours réseau — Vidéo 2 / 3"
    set_run(run, 14, False, ACCENT)

    # 23 AIDE MEMOIRE
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Aide-mémoire (screenshot)")
    add_timing(s, "Fiche")
    add_card(s, 0.5, 1.6, 6.0, 5.0, "Commandes switch",
             "show vlan brief\n"
             "show interfaces trunk\n"
             "show interfaces fa0/1 switchport\n\n"
             "switchport mode access\n"
             "switchport access vlan 10\n"
             "switchport mode trunk\n"
             "switchport trunk allowed vlan 10,20,30",
             LIGHT, size=14)
    add_card(s, 6.8, 1.6, 5.9, 5.0, "Commandes routeur",
             "interface g0/0.10\n"
             " encapsulation dot1Q 10\n"
             " ip address 10.10.10.1 255.255.255.0\n\n"
             "show ip route\n"
             "show ip interface brief\n\n"
             "ACL : deny invités → serveurs\nsur l’interface .20 in",
             RGBColor(0xE8, 0xF5, 0xE9), GREEN, 14)
    add_footer(s, 23)

    # 24 NOTES
    s = blank(prs)
    add_top_bar(s, "Notes formateur (ne pas montrer)")
    add_title(s, "Notes de tournage")
    add_bullets(s, [
        "Ouvrir Packet Tracer AVANT 34:00 pour ne pas perdre de temps",
        "Si switch PT sans « encapsulation dot1q » : passer la ligne",
        "Colorer les câbles PT (bleu VLAN10, vert 20, rouge 30, orange trunk)",
        "Montrer show vlan brief AVANT et APRÈS (effet wow)",
        "Ping inter-VLAN d’abord OK, PUIS ACL → le ping tombe : pédagogie",
        "Devoir VLAN 40 : ne pas trop détailler voice vlan à l’oral (teaser)",
        "CTA : « Screenshot show vlan brief en commentaire »",
        "Réponse devoir V1 à rappeler en intro si tu publies 7 j après : 10.0.4.0/23",
    ], top=1.65, size=16)
    add_footer(s, 24)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    stamp_meta(OUT)
    print(f"PowerPoint généré : {OUT} ({OUT.stat().st_size} octets) — {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
