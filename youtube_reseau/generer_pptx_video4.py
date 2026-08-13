#!/usr/bin/env python3
"""PowerPoint — Vidéo 4 : Wireshark pour débutants (1 h) — ABOU SAYABOU."""

from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Video4_Wireshark_Debutants_ABOU_SAYABOU.pptx"

NAVY = RGBColor(0x0D, 0x3B, 0x66)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
LIGHT = RGBColor(0xE8, 0xF1, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
RED = RGBColor(0xC6, 0x28, 0x28)
YELLOW_BG = RGBColor(0xFF, 0xF8, 0xE1)
RED_BG = RGBColor(0xFF, 0xEB, 0xEE)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
CODE_BG = RGBColor(0x1A, 0x23, 0x32)

W, H = Inches(13.333), Inches(7.5)
TOTAL = 25
FOOTER = "Vidéo 4 — Wireshark débutants  |  Débutant → Intermédiaire"


def set_run(run, size=20, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_top_bar(slide, title_right="Cours réseau — 1 h"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    run = bar.text_frame.paragraphs[0].add_run()
    run.text = "  ABOU SAYABOU  •  Réseau informatique"
    set_run(run, 12, True, WHITE)
    lab = slide.shapes.add_textbox(Inches(8.5), Inches(0.12), Inches(4.5), Inches(0.35))
    p = lab.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = title_right + "  "
    set_run(run, 12, False, LIGHT)


def add_footer(slide, page):
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), W, Inches(0.35))
    foot.fill.solid()
    foot.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    foot.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.18), Inches(10), Inches(0.3))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = FOOTER
    set_run(run, 10, False, GRAY)
    tb2 = slide.shapes.add_textbox(Inches(11.5), Inches(7.18), Inches(1.5), Inches(0.3))
    p = tb2.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page}/{TOTAL}"
    set_run(run, 10, False, GRAY)


def add_title(slide, text, top=0.7):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.7))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    set_run(run, 30, True, NAVY)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top + 0.62), Inches(1.8), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_subtitle(slide, text, top=1.42):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.4))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    set_run(run, 16, False, GRAY)


def add_bullets(slide, items, left=0.6, top=2.0, width=12, size=18):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "▸  " + item
        set_run(run, size, False, DARK)


def add_card(slide, left, top, width, height, title, body, fill=LIGHT, title_color=NAVY, size=13):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    run = tf.paragraphs[0].add_run()
    run.text = title
    set_run(run, 16, True, title_color)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    run = p2.add_run()
    run.text = body
    set_run(run, size, False, DARK)


def add_timing(slide, text):
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.68), Inches(2.4), Inches(0.4)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    set_run(run, 12, True, WHITE)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_code_card(slide, left, top, width, height, title, code):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CODE_BG
    card.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.1)
    run = tf.paragraphs[0].add_run()
    run.text = title
    set_run(run, 13, True, ACCENT, "Calibri")
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    run = p2.add_run()
    run.text = code
    set_run(run, 13, False, WHITE, "Consolas")


def stamp_meta(path: Path):
    now = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    buf = BytesIO(path.read_bytes())
    out_buf = BytesIO()
    with zipfile.ZipFile(buf, "r") as zin, zipfile.ZipFile(
        out_buf, "w", compression=zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "docProps/core.xml":
                data = f"""<?xml version='1.0' encoding='UTF-8' standalone='yes'?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <dc:title>Vidéo 4 — Wireshark pour débutants (cours 1h)</dc:title>
  <dc:subject>Cours réseau — capture de paquets, filtres, TCP handshake</dc:subject>
  <dc:creator>ABOU SAYABOU</dc:creator>
  <cp:keywords>Wireshark, pcap, filtres, TCP, DNS, YouTube</cp:keywords>
  <dc:description>Support de cours vidéo 1h — Wireshark débutants</dc:description>
  <cp:lastModifiedBy>ABOU SAYABOU</cp:lastModifiedBy>
  <cp:revision>1</cp:revision>
  <dcterms:created xsi:type="dcterms:W3CDTF">{now}</dcterms:created>
  <dcterms:modified xsi:type="dcterms:W3CDTF">{now}</dcterms:modified>
  <cp:category>Formation réseau</cp:category>
  <dc:language>fr-FR</dc:language>
</cp:coreProperties>""".encode("utf-8")
            elif item.filename == "docProps/app.xml":
                text = data.decode("utf-8")
                text = text.replace("<Company/>", "<Company>ABOU SAYABOU</Company>")
                text = text.replace("python-pptx", "ABOU SAYABOU")
                data = text.encode("utf-8")
            zout.writestr(item, data)
    path.write_bytes(out_buf.getvalue())


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 COVER
    s = blank(prs)
    add_bg(s, NAVY)
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.8), W, Inches(1.7))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(0x0A, 0x2F, 0x52)
    strip.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "COURS RÉSEAU  •  VIDÉO 4"
    set_run(run, 18, True, ACCENT)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.2))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Wireshark pour débutants :"
    set_run(run, 34, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    run = p2.add_run()
    run.text = "lire ses premiers paquets en 1 h"
    set_run(run, 28, True, LIGHT)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Capture  •  Interface  •  Filtres  •  TCP handshake  •  DNS / HTTPS"
    set_run(run, 16, False, RGBColor(0xB0, 0xC4, 0xDE))
    tb = s.shapes.add_textbox(Inches(0.8), Inches(6.15), Inches(11.5), Inches(0.8))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "ABOU SAYABOU"
    set_run(run, 22, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    run = p2.add_run()
    run.text = "Formateur cybersécurité & réseaux"
    set_run(run, 14, False, LIGHT)

    # 2 OBJECTIFS
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Objectif de la vidéo")
    add_timing(s, "0:00 – 2:00")
    add_card(s, 0.5, 1.9, 3.9, 4.4, "1. Capturer",
             "Lancer Wireshark, choisir la bonne interface, enregistrer un .pcapng propre.",
             LIGHT)
    add_card(s, 4.7, 1.9, 3.9, 4.4, "2. Filtrer",
             "Maîtriser 8 filtres utiles (ip, dns, tcp.port, http, tls…) sans se noyer.",
             LIGHT)
    add_card(s, 8.9, 1.9, 3.9, 4.4, "3. Lire",
             "Reconnaître un handshake TCP, une requête DNS, et ce que HTTPS cache.",
             LIGHT)
    add_footer(s, 2)

    # 3 PLAN
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Plan de la session (1 h)")
    add_timing(s, "Chapitres")
    plan = [
        ("01", "0:00", "Pourquoi Wireshark (et ce que ce n’est pas)"),
        ("02", "5:00", "Installation + cadre légal"),
        ("03", "10:00", "L’interface en 3 panneaux"),
        ("04", "18:00", "Première capture (lab guidé)"),
        ("05", "28:00", "Filtres d’affichage essentiels"),
        ("06", "38:00", "Lire une frame + TCP 3-way handshake"),
        ("07", "48:00", "DNS, HTTP, HTTPS — 3 cas concrets"),
        ("08", "56:00", "Astuces pro, récap, devoir"),
    ]
    y = 1.55
    for num, t, label in plan:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(0.7), Inches(0.52))
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = box.text_frame.paragraphs[0].add_run()
        run.text = num
        set_run(run, 14, True, WHITE)
        tbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.35), Inches(y), Inches(1.3), Inches(0.52))
        tbox.fill.solid()
        tbox.fill.fore_color.rgb = ACCENT
        tbox.line.fill.background()
        tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tbox.text_frame.paragraphs[0].add_run()
        run.text = t
        set_run(run, 13, True, WHITE)
        tb = s.shapes.add_textbox(Inches(2.9), Inches(y + 0.05), Inches(9.5), Inches(0.4))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = label
        set_run(run, 17, False, DARK)
        y += 0.6
    add_footer(s, 3)

    # 4 POURQUOI
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Pourquoi Wireshark ?")
    add_timing(s, "0:00 – 5:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Ça sert à",
             "• Voir le trafic réel (pas juste ping)\n"
             "• Prouver qu’un paquet part / arrive\n"
             "• Diagnostiquer DNS, TCP, lenteurs\n"
             "• Analyser un incident / malware\n"
             "• Comprendre TCP/IP « en vrai »\n"
             "• Préparer certifs / SOC / pentest\n\n"
             "C’est le microscope du réseau.",
             LIGHT, NAVY, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Ce que ce n’est PAS",
             "• Pas un scanner d’attaque\n"
             "• Pas magique sur Wi-Fi chiffré\n"
             "  (WPA2/3) sans clés\n"
             "• Pas lisible sur HTTPS\n"
             "  (contenu chiffré)\n"
             "• Pas légal partout sans mandat\n\n"
             "Outil de lecture, pas de piratage.",
             YELLOW_BG, ORANGE, 14)
    add_footer(s, 4)

    # 5 LEGAL
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Cadre légal & bonnes pratiques")
    add_timing(s, "5:00 – 10:00")
    add_bullets(s, [
        "Capture uniquement sur TON lab, TON PC, ou réseau AUTORISÉ (mandat écrit)",
        "Capturer le trafic d’autrui sans droit = risque pénal / RGPD",
        "En entreprise : politique de capture + accord RSSI / DSI",
        "Ne publie jamais un .pcap contenant mots de passe / tokens / données perso",
        "Pour YouTube / formation : génère du trafic lab (ping, DNS, HTTP clair)",
    ], top=1.7, size=17)
    tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.15))
    tip.fill.solid()
    tip.fill.fore_color.rgb = RED_BG
    tip.line.color.rgb = RED
    run = tip.text_frame.paragraphs[0].add_run()
    run.text = "Règle : si tu n’as pas le droit d’écouter ce câble / ce Wi-Fi, tu ne captures pas."
    set_run(run, 16, True, RED)
    add_footer(s, 5)

    # 6 INSTALL
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Installation rapide")
    add_timing(s, "Install")
    add_card(s, 0.5, 1.7, 4.0, 4.8, "Windows",
             "1. Télécharger wireshark.org\n"
             "2. Installer Npcap\n"
             "   (proposé à l’install)\n"
             "3. Lancer en admin si besoin\n"
             "4. Choisir l’interface\n\n"
             "Npcap = pilote de capture",
             LIGHT, NAVY, 14)
    add_card(s, 4.7, 1.7, 4.0, 4.8, "Linux",
             "sudo apt update\n"
             "sudo apt install wireshark\n\n"
             "Ajouter ton user au groupe\n"
             "wireshark (sans root) :\n\n"
             "sudo usermod -aG wireshark $USER\n"
             "# puis se reconnecter",
             GREEN_BG, GREEN, 13)
    add_card(s, 8.9, 1.7, 3.9, 4.8, "macOS",
             "Installer depuis\nwireshark.org\nou Homebrew :\n\n"
             "brew install --cask wireshark\n\n"
             "Autoriser les permissions\nde capture réseau.",
             YELLOW_BG, ORANGE, 14)
    add_footer(s, 6)

    # 7 INTERFACE
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "L’interface en 3 panneaux")
    add_timing(s, "10:00 – 18:00")
    add_card(s, 0.5, 1.7, 12.3, 1.35, "1. Liste des paquets (haut)",
             "Chaque ligne = 1 frame. Colonnes : No, Time, Source, Destination, Protocol, Length, Info.",
             LIGHT, NAVY, 15)
    add_card(s, 0.5, 3.2, 12.3, 1.35, "2. Détails du paquet (milieu)",
             "Arbre des couches : Frame → Ethernet → IP → TCP/UDP → Application (DNS, TLS…).",
             LIGHT, NAVY, 15)
    add_card(s, 0.5, 4.7, 12.3, 1.7, "3. Octets bruts (bas)",
             "Hexadécimal + ASCII. Utile pour preuve / extrait. Débutant : on regarde surtout 1 et 2.",
             LIGHT, NAVY, 15)
    add_footer(s, 7)

    # 8 COULEURS
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Couleurs des paquets (repères)")
    add_timing(s, "Repères")
    rows = [
        ("Vert clair", "TCP « normal »", "Trafic TCP usuel"),
        ("Bleu clair", "UDP", "DNS, etc."),
        ("Noir / rouge", "Erreur / problème", "Retransmissions, resets…"),
        ("Jaune", "À vérifier selon règles", "Souvent alertes / marquages"),
    ]
    y = 1.7
    for color, meaning, detail in rows:
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1.05))
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT
        c.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        tf = c.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = color + "  →  "
        set_run(run, 17, True, ACCENT)
        run = p.add_run()
        run.text = meaning + "  —  " + detail
        set_run(run, 16, False, DARK)
        y += 1.2
    add_footer(s, 8)

    # 9 PREMIERE CAPTURE
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Lab 1 — Première capture")
    add_timing(s, "18:00 – 28:00")
    add_bullets(s, [
        "Ouvre Wireshark → double-clic sur ton interface (Ethernet / Wi-Fi)",
        "Stop (carré rouge) après 10–20 secondes de trafic",
        "Génère du trafic contrôlé : ping 8.8.8.8  puis  nslookup google.com",
        "File → Save As → lab1_ping_dns.pcapng",
        "Si tu ne vois RIEN : mauvaise interface, droits, ou VPN qui encapsule",
    ], top=1.7, size=17)
    add_code_card(
        s, 0.5, 5.0, 12.3, 1.55,
        "Trafic à générer (terminal)",
        "ping -n 4 8.8.8.8                 # Windows\n"
        "nslookup google.com\n"
        "curl http://neverssl.com          # HTTP clair (lab)",
    )
    add_footer(s, 9)

    # 10 INTERFACE CHOIX
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Choisir la bonne interface")
    add_timing(s, "Astuce")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Indices",
             "• Regarde le « sparkline »\n"
             "  (petit graphique de trafic)\n"
             "• Ethernet vs Wi-Fi vs Loopback\n"
             "• VPN (tun/tap) = parfois\n"
             "  tout le trafic y passe\n"
             "• VirtualBox/VMware :\n"
             "  interfaces virtuelles\n\n"
             "Mauvaise interface =\npcap vide ou inutile.",
             LIGHT, NAVY, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Loopback",
             "127.0.0.1 / Adapter Loopback\n\n"
             "Utile pour trafic LOCAL\n(appli sur ta machine).\n\n"
             "Inutile pour voir Internet.\n\n"
             "Sous Windows : Npcap\nloopback souvent présent.",
             GREEN_BG, GREEN, 14)
    add_footer(s, 10)

    # 11 FILTRES INTRO
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Filtres : capture vs affichage")
    add_timing(s, "28:00 – 32:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Filtre de CAPTURE",
             "• Avant / pendant la capture\n"
             "• Réduit ce qui est enregistré\n"
             "• Syntaxe type BPF/libpcap\n"
             "  host 8.8.8.8\n"
             "  port 53\n\n"
             "Avantage : fichier plus léger\nInconvénient : tu peux rater\ndes preuves",
             YELLOW_BG, ORANGE, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Filtre d’AFFICHAGE",
             "• Après la capture (barre verte)\n"
             "• Cache / montre sans détruire\n"
             "• Syntaxe Wireshark :\n"
             "  dns\n"
             "  ip.addr == 8.8.8.8\n"
             "  tcp.port == 443\n\n"
             "Débutant : commence ICI.\nCapture large, filtre ensuite.",
             GREEN_BG, GREEN, 14)
    add_footer(s, 11)

    # 12 8 FILTRES
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "8 filtres d’affichage à connaître")
    add_timing(s, "32:00 – 38:00")
    filters = [
        ("dns", "Uniquement le DNS"),
        ("icmp", "Ping / traceroute ICMP"),
        ("arp", "Résolution IP → MAC"),
        ("http", "HTTP clair (pas HTTPS)"),
        ("tls  ou  tcp.port == 443", "Trafic HTTPS / TLS"),
        ("ip.addr == 192.168.1.10", "Tout vers/depuis cette IP"),
        ("tcp.port == 53 or udp.port == 53", "DNS (ports)"),
        ("tcp.flags.syn == 1", "Débuts de connexion TCP"),
    ]
    y = 1.5
    for f, meaning in filters:
        left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(5.8), Inches(0.58))
        left.fill.solid()
        left.fill.fore_color.rgb = CODE_BG
        left.line.fill.background()
        run = left.text_frame.paragraphs[0].add_run()
        run.text = "  " + f
        set_run(run, 13, True, WHITE, "Consolas")
        right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(y), Inches(6.3), Inches(0.58))
        right.fill.solid()
        right.fill.fore_color.rgb = LIGHT
        right.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        run = right.text_frame.paragraphs[0].add_run()
        run.text = "  " + meaning
        set_run(run, 14, False, DARK)
        y += 0.65
    add_footer(s, 12)

    # 13 OPERATEURS
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Combiner les filtres")
    add_timing(s, "Logique")
    add_code_card(
        s, 0.5, 1.6, 12.3, 5.0,
        "Exemples",
        "ip.addr == 10.10.10.10 and dns\n"
        "tcp.port == 443 and ip.src == 192.168.1.20\n"
        "http.request or http.response\n"
        "not arp and not stun          # nettoyer le bruit\n"
        "tcp.analysis.retransmission  # lenteurs / pertes\n\n"
        "and = ET    or = OU    not = SAUF\n"
        "La barre devient VERTE si le filtre est valide, ROUGE sinon.",
    )
    add_footer(s, 13)

    # 14 LIRE UNE FRAME
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Lire une frame (couche par couche)")
    add_timing(s, "38:00 – 42:00")
    layers = [
        ("Ethernet", "MAC source / destination"),
        ("IP", "IP src / dst, TTL"),
        ("TCP ou UDP", "Ports, flags, seq/ack"),
        ("Application", "DNS query, TLS Client Hello…"),
    ]
    y = 1.7
    for name, detail in layers:
        left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(3.5), Inches(1.0))
        left.fill.solid()
        left.fill.fore_color.rgb = NAVY
        left.line.fill.background()
        left.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = left.text_frame.paragraphs[0].add_run()
        run.text = "\n" + name
        set_run(run, 16, True, WHITE)
        right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(y), Inches(8.6), Inches(1.0))
        right.fill.solid()
        right.fill.fore_color.rgb = LIGHT
        right.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        run = right.text_frame.paragraphs[0].add_run()
        run.text = "  " + detail
        set_run(run, 17, False, DARK)
        y += 1.15
    add_footer(s, 14)

    # 15 TCP HANDSHAKE
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "TCP : le 3-way handshake")
    add_timing(s, "42:00 – 48:00")
    add_subtitle(s, "Comment une connexion TCP démarre — à voir dans Wireshark")
    steps = [
        ("1. SYN", "Client → Serveur", "« Je veux me connecter »"),
        ("2. SYN-ACK", "Serveur → Client", "« OK, je suis prêt »"),
        ("3. ACK", "Client → Serveur", "« Confirmé, on parle »"),
    ]
    x = 0.5
    for title, who, meaning in steps:
        add_card(s, x, 2.0, 4.0, 3.2, title, f"{who}\n\n{meaning}", LIGHT, NAVY, 15)
        x += 4.2
    add_code_card(
        s, 0.5, 5.4, 12.3, 1.2,
        "Filtre",
        "tcp.flags.syn == 1   ou   tcp.handshake    (selon version)",
    )
    add_footer(s, 15)

    # 16 FOLLOW STREAM
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Astuce pro : Follow TCP Stream")
    add_timing(s, "Astuce")
    add_bullets(s, [
        "Clic droit sur un paquet TCP → Follow → TCP Stream",
        "Wireshark reconstruit la conversation (client en rouge / serveur en bleu)",
        "Sur HTTP clair : tu vois GET, Host, User-Agent, HTML…",
        "Sur HTTPS : tu verras surtout du binaire chiffré (normal)",
        "Idéal pour comprendre une session sans sauter de paquet en paquet",
    ], top=1.7, size=17)
    tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.2))
    tip.fill.solid()
    tip.fill.fore_color.rgb = GREEN_BG
    tip.line.color.rgb = GREEN
    run = tip.text_frame.paragraphs[0].add_run()
    run.text = "C’est souvent LA fonction qui fait passer du débutant à l’intermédiaire."
    set_run(run, 16, True, GREEN)
    add_footer(s, 16)

    # 17 CAS DNS
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Cas 1 — Voir une résolution DNS")
    add_timing(s, "48:00 – 51:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Scénario lab",
             "1. Start capture\n"
             "2. nslookup google.com\n"
             "3. Stop\n"
             "4. Filtre : dns\n\n"
             "Cherche :\n"
             "• DNS query\n"
             "• DNS response\n"
             "• Adresse A / AAAA renvoyée",
             LIGHT, NAVY, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Ce que tu apprends",
             "• Port UDP 53 (souvent)\n"
             "• Qui est ton résolveur DNS\n"
             "• Temps de réponse\n"
             "• NXDOMAIN si le nom\n"
             "  n’existe pas\n\n"
             "Lien vidéo 3 :\n"
             "si ping IP OK et nom KO,\n"
             "le pcap DNS montre pourquoi.",
             GREEN_BG, GREEN, 14)
    add_footer(s, 17)

    # 18 CAS HTTP
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Cas 2 — HTTP clair (lab)")
    add_timing(s, "51:00 – 53:30")
    add_bullets(s, [
        "Va sur un site HTTP (ex. neverssl.com) pendant la capture",
        "Filtre : http",
        "Ouvre un GET → tu vois Host, chemin, éventuellement cookies",
        "Follow TCP Stream → conversation lisible",
        "Message clé : sur Internet réel, la plupart du web est en HTTPS",
    ], top=1.7, size=17)
    tip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.3))
    tip.fill.solid()
    tip.fill.fore_color.rgb = YELLOW_BG
    tip.line.color.rgb = ORANGE
    tf = tip.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = "Pédagogie : HTTP montre le principe. HTTPS protège le contenu — et c’est une bonne chose."
    set_run(run, 16, True, ORANGE)
    add_footer(s, 18)

    # 19 CAS HTTPS
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Cas 3 — HTTPS : que peut-on encore voir ?")
    add_timing(s, "53:30 – 56:00")
    add_card(s, 0.5, 1.7, 6.0, 4.8, "Visible",
             "• IP source / destination\n"
             "• Port 443\n"
             "• Handshake TLS (Client Hello)\n"
             "• SNI parfois (nom du site)\n"
             "• Taille / timing des paquets\n"
             "• Certificats (selon échange)\n\n"
             "Utile pour savoir QUOI\nparle à QUI.",
             LIGHT, NAVY, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.8, "Invisible (chiffré)",
             "• Mot de passe\n"
             "• Contenu de page\n"
             "• Cookies HTTPS\n"
             "• Corps des requêtes API\n\n"
             "Pour décrypter : lab avancé\n(clé SSLKEYLOGFILE) —\nhors scope débutant,\net uniquement sur TON lab.",
             RED_BG, RED, 14)
    add_footer(s, 19)

    # 20 STATISTIQUES
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Menus utiles (sans se perdre)")
    add_timing(s, "56:00")
    add_card(s, 0.5, 1.7, 4.0, 4.8, "Statistics",
             "• Protocol Hierarchy\n"
             "• Conversations\n"
             "• Endpoints\n"
             "• DNS / HTTP\n\n"
             "Vue macro avant\nde zoomer paquet\npar paquet.",
             LIGHT, NAVY, 14)
    add_card(s, 4.7, 1.7, 4.0, 4.8, "Analyze",
             "• Display Filters\n"
             "• Follow Stream\n"
             "• Expert Information\n"
             "  (warnings / errors)\n\n"
             "Expert Info = « où ça\nsaigne » dans le pcap.",
             GREEN_BG, GREEN, 14)
    add_card(s, 8.9, 1.7, 3.9, 4.8, "File",
             "• Save / Export\n"
             "• Export Objects\n"
             "  (HTTP fichiers)\n"
             "• Merge captures\n\n"
             "Toujours sauver le\n.pcapng du lab.",
             YELLOW_BG, ORANGE, 14)
    add_footer(s, 20)

    # 21 RECAP
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Récap — 6 phrases")
    add_timing(s, "58:00")
    recaps = [
        ("Capture", "Bonne interface + trafic contrôlé + .pcapng"),
        ("3 panneaux", "Liste → détails couches → hex"),
        ("Filtres", "Affichage d’abord : dns, icmp, ip.addr, tcp.port"),
        ("TCP", "SYN → SYN-ACK → ACK"),
        ("Follow", "TCP Stream pour lire une conversation"),
        ("HTTPS", "Métadonnées visibles, contenu chiffré"),
    ]
    y = 1.55
    for a, b in recaps:
        left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(3.0), Inches(0.8))
        left.fill.solid()
        left.fill.fore_color.rgb = NAVY
        left.line.fill.background()
        left.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = left.text_frame.paragraphs[0].add_run()
        run.text = "\n" + a
        set_run(run, 15, True, WHITE)
        right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.7), Inches(y), Inches(9.1), Inches(0.8))
        right.fill.solid()
        right.fill.fore_color.rgb = LIGHT
        right.line.color.rgb = RGBColor(0xBD, 0xD3, 0xE6)
        run = right.text_frame.paragraphs[0].add_run()
        run.text = "  " + b
        set_run(run, 15, False, DARK)
        y += 0.88
    add_footer(s, 21)

    # 22 AIDE MEMOIRE
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Aide-mémoire filtres (screenshot)")
    add_timing(s, "Fiche")
    add_code_card(
        s, 0.5, 1.55, 12.3, 5.05,
        "Copie / colle",
        "dns\n"
        "icmp\n"
        "arp\n"
        "http\n"
        "tls\n"
        "ip.addr == 192.168.1.10\n"
        "tcp.port == 443\n"
        "udp.port == 53\n"
        "tcp.flags.syn == 1\n"
        "tcp.analysis.retransmission\n"
        "ip.addr == 10.0.0.5 and dns",
    )
    add_footer(s, 22)

    # 23 DEVOIR + CTA
    s = blank(prs)
    add_top_bar(s)
    add_title(s, "Devoir + prochaine vidéo")
    add_timing(s, "59:00 – 60:00")
    add_card(s, 0.5, 1.7, 6.0, 4.6, "Devoir",
             "1. Capture 30 s sur ton PC\n"
             "2. Filtre dns + screenshot\n"
             "3. Trouve le handshake TCP\n"
             "   d’une connexion 443\n"
             "4. Follow TCP Stream sur\n"
             "   un HTTP lab (neverssl)\n\n"
             "Commente : « DNS OK »\nquand c’est fait.",
             YELLOW_BG, ORANGE, 14)
    add_card(s, 6.8, 1.7, 5.9, 4.6, "Suite playlist",
             "V1 IP / masque / GW\n"
             "V2 VLAN + lab\n"
             "V3 Diagnostic 7 étapes\n"
             "V4 Wireshark ← tu es ici\n\n"
             "V5 (prochaine) :\nFirewall PME en 1 h\n(règles, ports, NAT, logs)\n\n"
             "Abonne-toi pour la suite.",
             LIGHT, NAVY, 14)
    add_footer(s, 23)

    # 24 MERCI + NOTES combined - actually separate merci and notes
    s = blank(prs)
    add_bg(s, NAVY)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Merci !"
    set_run(run, 48, True, WHITE)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11.5), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Capturer  •  Filtrer  •  Lire  •  Follow Stream"
    set_run(run, 20, False, LIGHT)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "ABOU SAYABOU"
    set_run(run, 24, True, WHITE)
    p2 = tb.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run = p2.add_run()
    run.text = "Cours réseau — Vidéo 4 — Wireshark"
    set_run(run, 14, False, ACCENT)

    # 25 NOTES
    s = blank(prs)
    add_top_bar(s, "Notes formateur (ne pas montrer)")
    add_title(s, "Notes de tournage")
    add_bullets(s, [
        "Installer Wireshark AVANT le tournage ; tester Npcap / droits",
        "Préparer neverssl.com + nslookup + ping pour démos propres",
        "Zoomer sur la barre de filtre (vert/rouge) — très pédagogique",
        "Ne pas tenter de décrypter HTTPS en live (hors scope + risque)",
        "Montrer Follow TCP Stream sur HTTP, puis contraste HTTPS",
        "Rappeler V3 : les 4 ping + pcap DNS = combo gagnant",
        "Devoir : demander screenshot filtre dns (pas le pcap complet)",
        "Teaser V5 Firewall PME en fin de vidéo",
    ], top=1.65, size=16)
    add_footer(s, 25)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    stamp_meta(OUT)
    print(f"PowerPoint généré : {OUT} ({OUT.stat().st_size} octets) — {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
